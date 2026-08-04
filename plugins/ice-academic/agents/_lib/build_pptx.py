#!/usr/bin/env python3
"""
build_pptx.py — generic PPTX builder for sales agents.

Usage:
    python3 build_pptx.py spec.json out.pptx

FONT RULE (V02R03 · 2026.08.05) — ⚠ ข้อความเดิมที่ว่า "default = Tahoma ทุกข้อความ"
  **ยกเลิกแล้ว** (โค้ดไม่เคยอ่าน docstring นี้ · ตั้งแต่ V02R01 ฟอนต์มาจาก font_policy.RAILS)
  • ฟอนต์มาจาก **ราง** ที่ infer_rail() เดาจากชนิดเอกสาร (เอกชน/ราชการ)
  • สไลด์แน่น (>400 ตัวอักษร หรือ >8 บรรทัด หรือ ตาราง >40 ช่อง ในสไลด์ใดสไลด์หนึ่ง)
    → สลับทั้งเด็คเป็น 'Leelawadee' (ยอดวรรณยุกต์เตี้ยสุด = ไม่ชนเมื่อบีบบรรทัด)
    ปิด/บังคับด้วย spec["dense"] = false / true
  • spec["font_family"] ระบุเอง = เคารพ ไม่แทรกแซง (แต่ยังผ่านด่านนโยบาย V1/V2/V4/V5)

spec.json schema:
{
  "title": "Deck title",
  "subtitle": "Optional subtitle",
  "author": "Optional",
  "font_family": "Tahoma",   # optional, default = Tahoma (TH+EN balanced)
  "theme": {"primary": "#1F4E79", "accent": "#2E75B6"},   # optional
  "slides": [
    {"layout": "title", "title": "...", "subtitle": "..."},
    {"layout": "section", "title": "Section header"},
    {"layout": "bullets", "title": "...", "bullets": ["...", "..."]},
    {"layout": "two_column", "title": "...", "left": ["..."], "right": ["..."]},
    {"layout": "table", "title": "...", "headers": ["A","B"], "rows": [["x","y"]]},
    {"layout": "kpi", "title": "...", "kpis": [{"label":"ARR","value":"$1.2M","delta":"+18% YoY"}]},
    {"layout": "image", "title": "...", "image_path": "/path/to.png"},
    {"layout": "thanks", "title": "Thank you", "subtitle": "..."}
  ]
}
"""
import sys, os, json, subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls


def OxmlElement(tag):
    return parse_xml(f'<{tag} {nsdecls("a")}/>')


# ⭐ นโยบายฟอนต์มาจาก SSOT เดียว — ห้าม hard-code ชื่อฟอนต์ในไฟล์นี้ (V02R02)
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from font_policy import (RAILS, resolve_font_policy, infer_rail,
                         measure_slide_density, DENSE_FONT)   # noqa: E402


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── CHAR GUARD (PPTX Lesson #18) ────────────────────────────────────────────
# U+2192 "→" (และ arrow ญาติ) ทำให้ PowerPoint for Mac ปฏิเสธทั้งไฟล์ (Repair)
# ขณะที่ LibreOffice/qlmanage ปล่อยผ่าน (false-green). แทนด้วย ▸ (U+25B8) ที่
# เปิดได้ทุก engine + สื่อความ flow เดียวกัน. ตรวจซ้ำที่ deck_qa.py (safety net).
CHAR_REPLACEMENTS = {
    "→": "▸",  # → RIGHTWARDS ARROW       → ▸ BLACK RIGHT-POINTING SMALL TRIANGLE
    "⟶": "▸",  # ⟶ LONG RIGHTWARDS ARROW  → ▸
    "➜": "▸",  # ➜ HEAVY ROUND-TIPPED ARROW → ▸
    "➔": "▸",  # ➔ HEAVY WIDE-HEADED ARROW  → ▸
    "➙": "▸",  # ➙ HEAVY RIGHTWARDS ARROW   → ▸
}


def _sanitize_chars(obj, _stats=None):
    """Recursively replace PowerPoint-rejecting chars in any string within spec.
    Returns (sanitized_obj, replacement_count). Logs to stderr — never silent."""
    top = _stats is None
    if _stats is None:
        _stats = {"count": 0}
    if isinstance(obj, str):
        out = obj
        for bad, good in CHAR_REPLACEMENTS.items():
            if bad in out:
                _stats["count"] += out.count(bad)
                out = out.replace(bad, good)
        return (out, _stats["count"]) if top else out
    if isinstance(obj, dict):
        res = {k: _sanitize_chars(v, _stats) for k, v in obj.items()}
        return (res, _stats["count"]) if top else res
    if isinstance(obj, list):
        res = [_sanitize_chars(v, _stats) for v in obj]
        return (res, _stats["count"]) if top else res
    return (obj, _stats["count"]) if top else obj


def add_title_slide(prs, slide, theme):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    if len(s.placeholders) > 1:
        s.placeholders[1].text = slide.get("subtitle", "")


def add_section_slide(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")


def add_bullets_slide(prs, slide, theme):
    layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    body = s.placeholders[1].text_frame
    body.text = slide["bullets"][0] if slide.get("bullets") else ""
    for b in slide.get("bullets", [])[1:]:
        p = body.add_paragraph()
        p.text = b


def add_two_column(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    left_items = slide.get("left", [])
    right_items = slide.get("right", [])
    left = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)).text_frame
    right = s.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(4.5), Inches(5)).text_frame
    if left_items:
        left.text = left_items[0]
        for t in left_items[1:]:
            left.add_paragraph().text = t
    if right_items:
        right.text = right_items[0]
        for t in right_items[1:]:
            right.add_paragraph().text = t


def add_table(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    headers = slide.get("headers", [])
    rows = slide.get("rows", [])
    if not headers:
        return
    rows_n = len(rows) + 1
    cols_n = len(headers)
    tbl = s.shapes.add_table(rows_n, cols_n, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * rows_n)).table
    for j, h in enumerate(headers):
        tbl.cell(0, j).text = str(h)
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            tbl.cell(i, j).text = str(v)


def add_kpi(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    kpis = slide.get("kpis", [])
    n = max(1, len(kpis))
    box_w = 9 / n
    for i, k in enumerate(kpis):
        left = Inches(0.5 + i * box_w)
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(box_w - 0.2), Inches(2.5))
        tf = shape.text_frame
        tf.text = k.get("label", "")
        p = tf.add_paragraph()
        p.text = k.get("value", "")
        p.font.size = Pt(28)
        p.font.bold = True
        if k.get("delta"):
            p2 = tf.add_paragraph()
            p2.text = k["delta"]


def add_image(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    if slide.get("image_path"):
        s.shapes.add_picture(slide["image_path"], Inches(0.5), Inches(1.5), Inches(9), Inches(5))


# ── ⭐ D1 TRI-SLOT FONT BINDING (V02R01 — 2026.08.04) ────────────────────────
# บั๊กที่แก้: ไฟล์นี้ **ไม่เคย set ฟอนต์เลยสักบรรทัด** ตลอดอายุการใช้งาน — docstring
# บอก "default Tahoma" และ spec key `font_family` แต่โค้ดไม่เคยอ่านทั้งคู่ (dead doc)
# → output ได้ฟอนต์ default ของ python-pptx (Calibri) ซึ่งไม่มี glyph ไทยเลย
# ทำเป็น post-pass เดินทุก run แทนการแก้ layout function ทั้ง 8 ตัว (ปลอดภัยกว่า)
def _bind_font(tf, font):
    for p in tf.paragraphs:
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            for tag in ("a:latin", "a:ea", "a:cs"):          # ครบ 3 slot — D1
                for old in rPr.findall(qn(tag)):
                    rPr.remove(old)
                el = OxmlElement(tag)
                el.set("typeface", font)
                rPr.append(el)


def apply_font(prs, font):
    """เดินทุก shape/table cell ในทุก slide แล้วผูกฟอนต์ครบ 3 slot"""
    n = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                _bind_font(sh.text_frame, font); n += 1
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        _bind_font(cell.text_frame, font); n += 1
    return n


LAYOUTS = {
    "title": add_title_slide,
    "section": add_section_slide,
    "bullets": add_bullets_slide,
    "two_column": add_two_column,
    "table": add_table,
    "kpi": add_kpi,
    "image": add_image,
    "thanks": add_title_slide,
}


def build(spec_path, out_path):
    with open(spec_path) as f:
        spec = json.load(f)
    # CHAR GUARD (Lesson #18): auto-replace PowerPoint-rejecting chars (→ ▸) before build
    spec, _n = _sanitize_chars(spec)
    if _n:
        print(f"CHAR-GUARD: replaced {_n} arrow char(s) (U+2192/etc → ▸) — would have caused PowerPoint Repair", file=sys.stderr)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    theme = spec.get("theme", {})
    if not spec.get("slides"):
        spec["slides"] = [{"layout": "title", "title": spec.get("title", "Untitled"), "subtitle": spec.get("subtitle", "")}]
    elif spec["slides"][0].get("layout") != "title":
        spec["slides"].insert(0, {"layout": "title", "title": spec.get("title", "Untitled"), "subtitle": spec.get("subtitle", "")})
    for slide in spec["slides"]:
        layout = slide.get("layout", "bullets")
        LAYOUTS.get(layout, add_bullets_slide)(prs, slide, theme)

    # ⭐ D1 — ฟอนต์มาจากราง (§3.0) · spec override ได้เมื่อลูกค้า/แบรนด์บังคับ
    rail, _why = infer_rail(spec, out_path)
    print(f"📄 ราง: {rail}  ({_why})")
    if rail not in RAILS:
        sys.exit(f"rail ต้องเป็น {'|'.join(RAILS)} (ได้: {rail})")
    font = spec.get("font_family") or RAILS[rail]["font"]

    # ⭐ V02R03 (คำสั่ง user 2026.08.05) — สไลด์แน่น/ต้องบีบบรรทัด → DENSE_FONT ทั้งเด็ค
    #   เหตุผลเชิงตัวเลข: ยอดวรรณยุกต์ Leelawadee 0.737 em vs ฟอนต์ราง 0.924 em
    #   → เมื่อบีบ line-height ตัวที่ยอดสูงกว่าชนก่อน · PPTX ฝังฟอนต์ได้จึงไม่ห่วงเครื่องผู้รับ
    #   ปิดได้ด้วย spec["dense"] = false · บังคับเปิดด้วย spec["dense"] = true
    if not spec.get("font_family"):          # ผู้ใช้ระบุฟอนต์เองแล้ว = เคารพ ไม่แทรกแซง
        _forced = spec.get("dense")
        _dense, _dwhy = measure_slide_density(spec.get("slides", []))
        if _forced is True:
            _dense, _dwhy = True, "spec ระบุ dense=true"
        elif _forced is False:
            _dense, _dwhy = False, "spec ระบุ dense=false — ไม่สลับ"
        if _dense and font != DENSE_FONT:
            print(f"🎚 สไลด์แน่น → เปลี่ยนทั้งเด็คเป็น '{DENSE_FONT}': {_dwhy}")
            print(f"   (ยอดวรรณยุกต์ 0.737 em เทียบ '{font}' 0.924 em = ไม่ชนเมื่อบีบบรรทัด"
                  f" · แลกกับไทยเล็กกว่าละตินมากขึ้น · ปิดด้วย \"dense\": false)")
            font = DENSE_FONT
        elif not _dense:
            print(f"🎚 ความแน่น: {_dwhy} → ใช้ฟอนต์ราง")

    # ⭐ V02R02: ด่านนโยบายก่อนสร้างไฟล์ — ผิดนโยบาย = แก้ให้เป็นฟอนต์ราง + แจ้ง (ไม่ fail)
    font, _notices, _ = resolve_font_policy(font, rail, spec)
    for _n in _notices:
        print(_n)
    n = apply_font(prs, font)
    prs.save(out_path)
    print(f"OK: wrote {out_path} with {len(prs.slides)} slides · "
          f"font='{font}' (rail={rail}) ผูกครบ 3 slot ใน {n} text frame")

    # post-build gate — จุดตรวจเดียวกับทุกฟอร์แมต
    subprocess.run([sys.executable, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                                 "audit_fonts.py"), "--rail", rail,
                    *(["--allow-font", font] if spec.get("font_family") else []), out_path])


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("usage: build_pptx.py spec.json out.pptx", file=sys.stderr)
        sys.exit(2)
    build(sys.argv[1], sys.argv[2])
