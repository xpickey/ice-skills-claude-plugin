"""
office_readers.py — Shared Office/PDF/text extractors for iCE Compass.Next _lib.

Single source of truth for read_docx / read_pptx / read_xlsx / read_pdf /
read_md_or_txt, consolidating near-identical copies that lived in both
fact_checker.py and knowledge_sync.py (the fact_checker copies lacked the
parse-time try/except, so a malformed customer .docx could crash the QA gate).

Error contract (preserved from knowledge_sync, the safer prior version):
  Every reader returns a STRING. On failure it returns a sentinel string that
  starts with "[ERROR:". Callers that care (e.g. knowledge_sync.sync_agent_docs)
  detect failure via content.startswith("[ERROR:"). This in-band contract is kept
  deliberately so existing callers keep working unchanged; fact_checker now gets
  the same graceful-failure behavior for free instead of crashing on a bad file.

V01R01 | 2026.06.08
"""

from pathlib import Path

ERROR_PREFIX = "[ERROR:"


def read_docx(path: Path) -> str:
    """Extract text from .docx (paragraphs + table cells)."""
    path = Path(path)
    try:
        from docx import Document
    except ImportError:
        return f"[ERROR: python-docx not installed — cannot read {path.name}]"
    try:
        d = Document(path)
        parts = [p.text for p in d.paragraphs if p.text.strip()]
        for tbl in d.tables:
            for row in tbl.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR: docx read failed — {e}]"


def read_pptx(path: Path) -> str:
    """Extract text from .pptx (per-slide shapes)."""
    path = Path(path)
    try:
        from pptx import Presentation
    except ImportError:
        return f"[ERROR: python-pptx not installed — cannot read {path.name}]"
    try:
        p = Presentation(path)
        parts = []
        for i, slide in enumerate(p.slides, 1):
            parts.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    if shape.text_frame.text.strip():
                        parts.append(shape.text_frame.text)
                elif hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR: pptx read failed — {e}]"


def read_xlsx(path: Path) -> str:
    """Extract text + numbers from .xlsx (all sheets, data_only)."""
    path = Path(path)
    try:
        from openpyxl import load_workbook
    except ImportError:
        return f"[ERROR: openpyxl not installed — cannot read {path.name}]"
    try:
        wb = load_workbook(path, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(v) if v is not None else "" for v in row)
                if row_text.strip(" |"):
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR: xlsx read failed — {e}]"


def read_pdf(path: Path) -> str:
    """Extract text from .pdf via pdftotext (poppler) if available."""
    path = Path(path)
    import subprocess
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", str(path), "-"],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0:
            return result.stdout
        return f"[ERROR: pdftotext returned {result.returncode}: {result.stderr}]"
    except FileNotFoundError:
        return f"[ERROR: pdftotext not installed — cannot read {path.name}. Install via 'brew install poppler']"
    except Exception as e:
        return f"[ERROR: pdf read failed — {e}]"


def read_md_or_txt(path: Path) -> str:
    """Read a plain text or markdown file as utf-8."""
    path = Path(path)
    try:
        return path.read_text(encoding="utf-8")
    except Exception as e:
        return f"[ERROR: text read failed — {e}]"
