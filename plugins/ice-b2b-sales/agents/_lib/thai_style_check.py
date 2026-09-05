#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""ตัวตรวจภาษาแปลและสำนวน AI ในข้อความไทย — thai_style_check.py (V01R01 · 2026.09.05)

ตรวจด้วยเครื่องก่อนถึงผู้ตรวจคุณภาพ ตามกติกาใน skill ice-writing-register (ส่วนที่ 2 และ 3) และหลักการเลี่ยงภาษาแปล
ของสำนักงานราชบัณฑิตยสภา ("การใช้ภาษาไทยในบทแปล") — ใช้ PyThaiNLP ตัดคำเพื่อวัดความยาวประโยคและความหนาแน่นของคำ

สิ่งที่ตรวจ (แต่ละข้อรายงานตำแหน่งและข้อความ):
  T1 สำนวนแปลสำเร็จรูป      "เป็นที่น่าสังเกตว่า" "มีบทบาทสำคัญอย่างยิ่ง" "ในโลกปัจจุบันที่" "ไม่เพียงแต่…แต่ยัง" ฯลฯ
  T2 กรรมวาจกแบบอังกฤษ      "ถูก…โดย" ที่ไม่ใช่เรื่องร้าย (ภาษาไทยใช้ "ถูก" กับเรื่องไม่ดี)
  T3 คำเติมไร้ความหมาย        "ทำการ" "มีการ" "ในส่วนของ" "ทางด้าน" นำหน้ากริยา/นาม
  T4 คำเชื่อมซ้อน             "ซึ่ง" หรือ "โดย" เกิน 2 ครั้งในประโยคเดียว
  T5 คำติดปาก AI ไทย          "ยิ่งไปกว่านั้น" "ในท้ายที่สุด" "อย่างก้าวกระโดด" "พลิกโฉม" "เป็นที่ทราบกันดี" · ประโยชน์สามท่อนลอย
  T6 คำติดปาก AI อังกฤษ       delve leverage utilize robust seamless comprehensive pivotal foster navigate showcase underscore streamline holistic
  T7 คำต้องห้ามงานธุรกิจ       zero-risk best-in-class world-class · คำละคร องก์ ฉาก ตัวละคร
  T8 ประโยคยาวเท่ากันหมด      ความยาวประโยค (คำ) ไม่แกว่งเลย = ขาดจังหวะของคนเขียน (เตือน)
  T9 คำที่อาจประดิษฐ์เอง       คำไทยที่ตัดคำแล้วได้ชิ้นส่วนที่ไม่อยู่ในพจนานุกรม ติดกันตั้งแต่ 2 ชิ้น (เตือนให้ตรวจกับ glossary)

วิธีใช้:  python3 ~/.claude/agents/_lib/thai_style_check.py FILE.(md|txt|docx|pptx) [--json] [--register business|academic]
ผลลัพธ์:  รายการที่พบ + สรุป · exit 0 = ไม่พบข้อระงับ · exit 2 = พบข้อระงับ (T1/T5/T7) ต้องแก้ก่อนส่งตรวจ
"""
import argparse
import json
import re
import statistics
import sys

try:
    from pythainlp.tokenize import word_tokenize, sent_tokenize
    from pythainlp.corpus import thai_words
    HAVE_NLP = True
except ImportError:
    HAVE_NLP = False

T1 = ["เป็นที่น่าสังเกตว่า", "มีบทบาทสำคัญอย่างยิ่ง", "ในโลกปัจจุบันที่", "ไม่เพียงแต่", "อย่างไรก็ตาม,", "ในบริบทของ", "เป็นสิ่งสำคัญที่จะต้อง", "เป็นที่ยอมรับกันโดยทั่วไป"]
T3 = re.compile(r"(ทำการ|มีการ)(?=[ก-๙])|ในส่วนของ|ทางด้าน(?=[ก-๙])")
T5 = ["ยิ่งไปกว่านั้น", "ในท้ายที่สุด", "อย่างก้าวกระโดด", "พลิกโฉม", "เป็นที่ทราบกันดี", "ทั้งหมดนี้สังเคราะห์ได้ว่า", "ปลดล็อกศักยภาพ", "ขับเคลื่อนอย่างยั่งยืน"]
T5_TRIAD = re.compile(r"(ทันสมัย|คล่องตัว|มีประสิทธิภาพ|ยั่งยืน|โปร่งใส|รวดเร็ว|แม่นยำ)[ ,]+(ทันสมัย|คล่องตัว|มีประสิทธิภาพ|ยั่งยืน|โปร่งใส|รวดเร็ว|แม่นยำ)[ ,]+(และ)?(ทันสมัย|คล่องตัว|มีประสิทธิภาพ|ยั่งยืน|โปร่งใส|รวดเร็ว|แม่นยำ)")
T6 = re.compile(r"\b(delve|leverag\w*|utiliz\w*|robust|seamless\w*|comprehensive|pivotal|foster\w*|navigat\w*|showcas\w*|underscor\w*|streamlin\w*|holistic|cutting-edge|tapestry|testament)\b", re.I)
T7 = re.compile(r"zero-risk|best-in-class|world-class|\bองก์\b|ฉากที่|ตัวละคร", re.I)
PASSIVE = re.compile(r"ถูก([ก-๙]{2,12})โดย")
BAD_PASSIVE_OK = ("ยกเลิก", "ปฏิเสธ", "ระงับ", "ตัด", "ลงโทษ", "ปรับ", "ฟ้อง")


def read_text(path):
    if path.endswith(".docx"):
        from docx import Document
        d = Document(path)
        parts = [p.text for p in d.paragraphs]
        for t in d.tables:
            for row in t.rows:
                parts += [c.text for c in row.cells]
        return "\n".join(parts)
    if path.endswith(".pptx"):
        from pptx import Presentation
        out = []
        for i, s in enumerate(Presentation(path).slides, 1):
            for sh in s.shapes:
                if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                    out.append(f"[หน้า {i}] " + sh.text_frame.text)
        return "\n".join(out)
    return open(path, encoding="utf-8", errors="ignore").read()


def sentences(text):
    if HAVE_NLP:
        try:
            return [s for s in sent_tokenize(text, engine="crfcut") if s.strip()]
        except Exception:
            pass
    return [s for s in re.split(r"[\n。.!?]+|(?<=[ก-๙]) (?=[ก-๙])", text) if s.strip()]


def check(text, register="business"):
    finds = []
    for kw in T1:
        for m in re.finditer(re.escape(kw), text):
            finds.append(("T1", "block", kw, ctx(text, m.start())))
    for m in PASSIVE.finditer(text):
        if not any(k in m.group(1) for k in BAD_PASSIVE_OK):
            finds.append(("T2", "warn", m.group(0), ctx(text, m.start())))
    for m in T3.finditer(text):
        finds.append(("T3", "warn", m.group(0), ctx(text, m.start())))
    for s in sentences(text):
        if s.count("ซึ่ง") + s.count("โดย") > 2:
            finds.append(("T4", "warn", f"ซึ่ง/โดย {s.count('ซึ่ง') + s.count('โดย')} ครั้ง", s[:80]))
    for kw in T5:
        for m in re.finditer(re.escape(kw), text):
            finds.append(("T5", "block", kw, ctx(text, m.start())))
    for m in T5_TRIAD.finditer(text):
        finds.append(("T5", "block", "ประโยชน์สามท่อนลอย: " + m.group(0), ctx(text, m.start())))
    for m in T6.finditer(text):
        finds.append(("T6", "warn", m.group(0), ctx(text, m.start())))
    if register == "business":
        for m in T7.finditer(text):
            finds.append(("T7", "block", m.group(0), ctx(text, m.start())))
    lens = []
    if HAVE_NLP:
        for s in sentences(text):
            w = [t for t in word_tokenize(s, keep_whitespace=False) if t.strip()]
            if len(w) >= 4:
                lens.append(len(w))
        if len(lens) >= 6 and statistics.pstdev(lens) < 2.5:
            finds.append(("T8", "warn", f"ประโยคยาวใกล้เคียงกันหมด (เฉลี่ย {statistics.mean(lens):.0f} คำ ส่วนเบี่ยงเบน {statistics.pstdev(lens):.1f})", "ทั้งเอกสาร"))
        vocab = thai_words()
        unknown_runs = set()
        for s in sentences(text):
            toks = [t for t in word_tokenize(s, keep_whitespace=False) if re.fullmatch(r"[ก-๙]+", t)]
            run = []
            for t in toks:
                if t not in vocab and len(t) >= 2:
                    run.append(t)
                else:
                    if len(run) >= 2:
                        unknown_runs.add("".join(run))
                    run = []
            if len(run) >= 2:
                unknown_runs.add("".join(run))
        for u in sorted(unknown_runs)[:20]:
            finds.append(("T9", "warn", u, "คำที่พจนานุกรมไม่รู้จักติดกัน — ตรวจกับ glossary หรือเอกสารลูกค้าว่าเป็นคำที่ใช้จริง"))
    return finds


def ctx(text, i, n=45):
    return text[max(0, i - n):i + n].replace("\n", " ")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("file")
    ap.add_argument("--register", choices=["business", "academic", "general"], default="business")
    ap.add_argument("--json", action="store_true")
    a = ap.parse_args()
    text = read_text(a.file)
    finds = check(text, a.register)
    blocks = [f for f in finds if f[1] == "block"]
    if a.json:
        print(json.dumps({"file": a.file, "blocks": len(blocks), "warns": len(finds) - len(blocks), "findings": finds}, ensure_ascii=False))
    else:
        print(f"== {a.file} · ชนิดงาน {a.register} · PyThaiNLP {'พร้อม' if HAVE_NLP else 'ไม่มี (ตรวจเฉพาะรูปแบบคำ)'} ==")
        for code, lvl, hit, c in finds:
            print(f"  {code} [{'ต้องแก้' if lvl == 'block' else 'เตือน'}] {hit}  …{c}…")
        print(f"  ผล: {'ต้องแก้ ' + str(len(blocks)) + ' จุด' if blocks else 'ไม่พบข้อต้องแก้'} · เตือน {len(finds) - len(blocks)} จุด")
    sys.exit(2 if blocks else 0)


if __name__ == "__main__":
    main()
