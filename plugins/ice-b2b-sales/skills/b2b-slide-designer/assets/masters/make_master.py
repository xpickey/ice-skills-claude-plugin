#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
make_master.py — สร้างแม่แบบ iCE-Propose_Master.pptx (16:9 · 13.333 × 7.5 นิ้ว) จากศูนย์
V01R01 | 2026.09.05 | Wave B ของแผนแก้ทีม agent (ต้นตอ: builder วาดทุกอย่างจาก shape primitive
เฉลี่ย 47 shape ต่อหน้า วัตถุทับกัน ฟอนต์ไม่คงที่ หน้าตาไม่เหมือน CI)

ทำอะไร
  1. เริ่มจาก Presentation() ของ python-pptx (มี master 1 ตัว + layout 11 ตัวติดมา)
  2. แก้ theme1.xml: สี theme จาก tokens.json ของ iCE Design System และฟอนต์ theme ทั้งสามช่อง
     (latin / ea / cs) = ฟอนต์รางเอกชนจาก font_policy.RAILS (ห้าม hard-code ชื่อฟอนต์ในไฟล์นี้)
  3. เขียน slideMaster ใหม่ทั้งตัว (พื้นหลังไล่เฉดอ่อน · โลโก้ · เส้นทอง · footer · เลขหน้า · txStyles)
  4. เขียน layout 10 แบบทับ layout เดิม 10 ตัวแรก แล้วลบตัวที่เหลือ:
       cover · divider · action-title-body · two-column · three-card · table · timeline · closing · appendix · blank
     ทุก layout ใช้ placeholder จริง (title / body / pic / tbl / ftr / sldNum) ไม่ใช่ textbox ลอย
  5. บันทึกเป็น iCE-Propose_Master.pptx ในโฟลเดอร์เดียวกับสคริปต์นี้

วิธีใช้
  python3 make_master.py            → เขียน iCE-Propose_Master.pptx ข้างสคริปต์
  python3 make_master.py OUT.pptx   → เขียนไปที่ path ที่ระบุ

ตรวจหลังสร้าง (ทำทุกครั้ง)
  python3 ~/.claude/agents/_lib/validate_pptx_structure.py iCE-Propose_Master.pptx
  bash    ~/.claude/agents/_lib/render_pdf.sh iCE-Propose_Master.pptx _review

ที่มาของค่าออกแบบ
  สี/ไล่เฉด/ทอง   : /Users/xpickey/Documents/Claude/iCE-Design-System/tokens/tokens.json
  หน้าตาต้นแบบ     : /Users/xpickey/Documents/Claude/iCE-Design-System/presentation/{presentation,presentation-light}.html
                     (ปก/คั่น/ปิด = พื้นเข้มไล่เฉดน้ำเงิน→ฟ้า + เส้นทองบาง · หน้าเนื้อหา = พื้นขาว โลโก้สี)
  โลโก้             : /Users/xpickey/Documents/Claude/iCE-Design-System/assets/logos/ice-logo-full.png (สี) · ice-logo-white.png (ขาว)
  กฎออกแบบ         : b2b-slide-designer/references/pptx-design-doctrine.md (กฎ 6 ข้อ)
  ฟอนต์            : ~/.claude/agents/_lib/font_policy.py (ราง private) + ice-doc-builder D1/D3
  ขนาดตัวอักษร      : D3 — ไทย = อังกฤษ ไม่บวก pt · เนื้อหา ≥16pt · หัวเรื่อง ≥24pt · เลี่ยง Bold ที่เนื้อหา
"""
import copy
import json
import os
import sys
from lxml import etree

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.shapes.picture import CT_Picture
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

sys.path.insert(0, os.path.expanduser("~/.claude/agents/_lib"))
from font_policy import RAILS  # noqa: E402  (SSOT ของฟอนต์ — ห้าม hard-code)

HERE = os.path.dirname(os.path.abspath(__file__))
DS = "/Users/xpickey/Documents/Claude/iCE-Design-System"
TOKENS = os.path.join(DS, "tokens", "tokens.json")
LOGO_COLOR = os.path.join(DS, "assets", "logos", "ice-logo-full.png")
LOGO_WHITE = os.path.join(DS, "assets", "logos", "ice-logo-white.png")
LOGO_RATIO = 2001 / 997          # กว้าง:สูง ของไฟล์โลโก้จริง

FONT = RAILS["private"]["font"]  # รางเอกชน — ค่าเดียวที่ใช้ทุกช่อง latin/ea/cs

EMU = 914400
W_IN, H_IN = 13.333, 7.5
W, H = 12192000, 6858000         # 16:9 มาตรฐาน (EMU)

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
NSDECL = ' '.join(f'xmlns:{k}="{v}"' for k, v in NS.items())


def emu(x_in):
    return int(round(x_in * EMU))


# ─────────────────────────────────────────────────────────────────────────────
# สีจาก tokens.json (single source) — คีย์ที่ใช้จริงในแม่แบบ
# ─────────────────────────────────────────────────────────────────────────────
def load_colors():
    with open(TOKENS, encoding="utf-8") as f:
        t = json.load(f)
    c = t["color"]
    hx = lambda v: v.lstrip("#").upper()
    col = {
        "navy": hx(c["brand"]["navy"]["value"]),
        "teal": hx(c["brand"]["teal"]["value"]),
        "navyDeep": hx(c["brand"]["navyDeep"]["value"]),
        "tealBright": hx(c["brand"]["tealBright"]["value"]),
        "ink": hx(c["brand"]["ink"]["value"]),
        "white": hx(c["brand"]["white"]["value"]),
        "gold": hx(c["gold"]["base"]["value"]),
        "goldLight": hx(c["gold"]["light"]["value"]),
        "goldDeep": hx(c["gold"]["deep"]["value"]),
        "success": hx(c["status"]["success"]["value"]),
        "warning": hx(c["status"]["warning"]["value"]),
        "danger": hx(c["status"]["danger"]["value"]),
    }
    # จุดสีของไล่เฉดที่ประกาศไว้ใน tokens (อ่านจากสตริง CSS ตรง ๆ เพื่อไม่ต้องคัดลอกค่า)
    import re
    grad = lambda k: [s.upper() for s in re.findall(r"#([0-9A-Fa-f]{6})", t["gradient"][k]["value"])]
    col["grad_dark"] = grad("dark")      # 0E2A47 → 15375C → 1E5E78
    col["grad_light"] = grad("light")    # FFFFFF → F2F7FB → EAF1F7
    col["grad_brand"] = grad("brand")    # 1E66A4 → 41A8B5
    return col


C = load_colors()
# สีอ่อนสำหรับการ์ด/เส้นขอบ (จาก iCE-Propose designer brief §2.3 — tint ของสี CI ไม่ใช่สีใหม่)
CARD_BG = "F4F6F8"
CARD_LINE = "A8C7E0"
TITLE_INK = C["navyDeep"]
BODY_INK = C["ink"]
MUTED = "8C8C8C"


# ─────────────────────────────────────────────────────────────────────────────
# ตัวช่วยประกอบ XML
# ─────────────────────────────────────────────────────────────────────────────
def xfrm(x, y, w, h):
    return (f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/>'
            f'<a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>')


def rpr_font():
    return (f'<a:latin typeface="{FONT}"/><a:ea typeface="{FONT}"/><a:cs typeface="{FONT}"/>')


def solid(hexcol, alpha=None):
    a = f'<a:alpha val="{int(alpha * 1000)}"/>' if alpha is not None else ''
    return f'<a:solidFill><a:srgbClr val="{hexcol}">{a}</a:srgbClr></a:solidFill>'


def grad_fill(stops, angle_deg=135):
    """stops = [(pos_percent, hex, alpha%|None), ...] · angle ตามแบบ CSS (135 = ซ้ายบน→ขวาล่าง)"""
    gs = ''
    for pos, hexcol, alpha in stops:
        a = f'<a:alpha val="{int(alpha * 1000)}"/>' if alpha is not None else ''
        gs += f'<a:gs pos="{int(pos * 1000)}"><a:srgbClr val="{hexcol}">{a}</a:srgbClr></a:gs>'
    # CSS 135deg → OOXML lin ang: 0 = ซ้าย→ขวา หมุนตามเข็ม · 135deg CSS ≈ 45° OOXML (2700000 = 45°)
    ang = {135: 2700000, 90: 0, 180: 5400000, 120: 1800000}.get(angle_deg, 2700000)
    return f'<a:gradFill rotWithShape="1"><a:gsLst>{gs}</a:gsLst><a:lin ang="{ang}" scaled="0"/></a:gradFill>'


def ln(hexcol, w_pt=0.75, alpha=None):
    return f'<a:ln w="{int(w_pt * 12700)}">{solid(hexcol, alpha)}</a:ln>'


def no_ln():
    return '<a:ln><a:noFill/></a:ln>'


_ids = {}


def next_id(part_key):
    _ids[part_key] = _ids.get(part_key, 1) + 1
    return _ids[part_key]


def sp_rect(key, name, x, y, w, h, fill_xml, line_xml=no_ln(), prst="rect", text_xml=None, extra_bodypr=''):
    """รูปทรงตกแต่งของ master/layout (ไม่ใช่ placeholder)"""
    tx = (f'<p:txBody><a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" anchor="ctr" rtlCol="0"{extra_bodypr}/>'
          f'<a:lstStyle/>{text_xml}</p:txBody>') if text_xml else ''
    return (f'<p:sp><p:nvSpPr><p:cNvPr id="{next_id(key)}" name="{name}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
            f'<p:spPr>{xfrm(x, y, w, h)}<a:prstGeom prst="{prst}"><a:avLst/></a:prstGeom>{fill_xml}{line_xml}</p:spPr>'
            f'{tx}</p:sp>')


def cxn_line(key, name, x1, y1, x2, y2, hexcol, w_pt=0.75, alpha=None, grad=None):
    """เส้นตรง (connector) — ใช้ทำเส้นทองบางและเส้นแกนเวลา"""
    x, y = min(x1, x2), min(y1, y2)
    w, h = abs(x2 - x1), abs(y2 - y1)
    fill = grad if grad else solid(hexcol, alpha)
    return (f'<p:cxnSp><p:nvCxnSpPr><p:cNvPr id="{next_id(key)}" name="{name}"/><p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
            f'<p:spPr>{xfrm(x, y, w, h)}<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
            f'<a:ln w="{int(w_pt * 12700)}" cap="rnd">{fill}<a:round/></a:ln></p:spPr></p:cxnSp>')


def para(text, sz, color, bold=False, align="l", spc_after=0):
    b = ' b="1"' if bold else ''
    return (f'<a:p><a:pPr algn="{align}"><a:spcAft><a:spcPts val="{spc_after * 100}"/></a:spcAft></a:pPr>'
            f'<a:r><a:rPr lang="th-TH" sz="{sz * 100}"{b} dirty="0">{solid(color)}{rpr_font()}</a:rPr>'
            f'<a:t>{text}</a:t></a:r></a:p>')


def ph(key, name, ph_type, idx, x, y, w, h, sz, color, bold=False, align="l", anchor="t",
       prompt="", lvl_extra="", bullets=False, sz2=None, autofit=True, cap=None):
    """placeholder จริง — ทุกช่องข้อความในแม่แบบใช้ตัวนี้"""
    idx_attr = f' idx="{idx}"' if idx is not None else ''
    type_attr = f' type="{ph_type}"' if ph_type else ''
    sz_attr = ' sz="quarter"' if ph_type in ("ftr", "sldNum", "dt") else ''
    b = ' b="1"' if bold else ''
    fit = '<a:normAutofit/>' if autofit else ''
    if bullets:
        lvl1 = (f'<a:lvl1pPr marL="228600" indent="-228600" algn="{align}">'
                f'<a:spcBef><a:spcPts val="600"/></a:spcBef><a:buClr><a:srgbClr val="{C["teal"]}"/></a:buClr>'
                f'<a:buFont typeface="{FONT}"/><a:buChar char="•"/>'
                f'<a:defRPr sz="{sz * 100}"{b}>{solid(color)}{rpr_font()}</a:defRPr></a:lvl1pPr>'
                f'<a:lvl2pPr marL="457200" indent="-228600" algn="{align}">'
                f'<a:spcBef><a:spcPts val="300"/></a:spcBef><a:buClr><a:srgbClr val="{C["tealBright"]}"/></a:buClr>'
                f'<a:buFont typeface="{FONT}"/><a:buChar char="–"/>'
                f'<a:defRPr sz="{(sz2 or sz - 2) * 100}">{solid(color)}{rpr_font()}</a:defRPr></a:lvl2pPr>')
    else:
        lvl1 = (f'<a:lvl1pPr marL="0" indent="0" algn="{align}"><a:buNone/>'
                f'<a:defRPr sz="{sz * 100}"{b}>{solid(color)}{rpr_font()}</a:defRPr></a:lvl1pPr>'
                f'<a:lvl2pPr marL="0" indent="0" algn="{align}"><a:buNone/>'
                f'<a:defRPr sz="{(sz2 or sz) * 100}">{solid(color)}{rpr_font()}</a:defRPr></a:lvl2pPr>')
    cap_attr = f' cap="{cap}"' if cap else ''
    body = (f'<p:txBody><a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" rtlCol="0" anchor="{anchor}">{fit}</a:bodyPr>'
            f'<a:lstStyle>{lvl1}{lvl_extra}</a:lstStyle>'
            f'<a:p><a:r><a:rPr lang="th-TH"{cap_attr} dirty="0"/><a:t>{prompt}</a:t></a:r></a:p></p:txBody>')
    return (f'<p:sp><p:nvSpPr><p:cNvPr id="{next_id(key)}" name="{name}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph{type_attr}{sz_attr}{idx_attr}/></p:nvPr></p:nvSpPr>'
            f'<p:spPr>{xfrm(x, y, w, h)}<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>{body}</p:sp>')


def ph_pic(key, name, idx, x, y, w, h, prompt="ภาพ / icon"):
    return (f'<p:sp><p:nvSpPr><p:cNvPr id="{next_id(key)}" name="{name}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="pic" idx="{idx}"/></p:nvPr></p:nvSpPr>'
            f'<p:spPr>{xfrm(x, y, w, h)}<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
            f'<p:txBody><a:bodyPr anchor="ctr"/><a:lstStyle><a:lvl1pPr algn="ctr"><a:buNone/>'
            f'<a:defRPr sz="1200">{solid(MUTED)}{rpr_font()}</a:defRPr></a:lvl1pPr></a:lstStyle>'
            f'<a:p><a:r><a:rPr lang="th-TH" dirty="0"/><a:t>{prompt}</a:t></a:r></a:p></p:txBody></p:sp>')


def ph_tbl(key, name, idx, x, y, w, h):
    return (f'<p:sp><p:nvSpPr><p:cNvPr id="{next_id(key)}" name="{name}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="tbl" idx="{idx}"/></p:nvPr></p:nvSpPr>'
            f'<p:spPr>{xfrm(x, y, w, h)}</p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle><a:lvl1pPr algn="ctr"><a:buNone/>'
            f'<a:defRPr sz="1400">{solid(MUTED)}{rpr_font()}</a:defRPr></a:lvl1pPr></a:lstStyle>'
            f'<a:p><a:r><a:rPr lang="th-TH" dirty="0"/><a:t>ตาราง</a:t></a:r></a:p></p:txBody></p:sp>')


def footer_phs(key, dark=False):
    """footer (idx 11) + เลขหน้า (idx 12) — ตำแหน่งเดียวกันทุก layout"""
    col = "D9E4EE" if dark else MUTED
    ftr = ph(key, "Footer Placeholder", "ftr", 11, 0.6, 6.86, 9.0, 0.3, 10, col, anchor="ctr",
             prompt="iCE Consulting · เอกสารลับ", autofit=False)
    num = (f'<p:sp><p:nvSpPr><p:cNvPr id="{next_id(key)}" name="Slide Number Placeholder"/>'
           f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldNum" sz="quarter" idx="12"/></p:nvPr></p:nvSpPr>'
           f'<p:spPr>{xfrm(11.73, 6.86, 1.0, 0.3)}<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
           f'<p:txBody><a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" rtlCol="0" anchor="ctr"/>'
           f'<a:lstStyle><a:lvl1pPr algn="r"><a:buNone/><a:defRPr sz="1000">{solid(col)}{rpr_font()}</a:defRPr></a:lvl1pPr></a:lstStyle>'
           f'<a:p><a:fld id="{{B6F15528-21DE-4FAA-801E-634DDDAF4B2B}}" type="slidenum"><a:rPr lang="en-US" dirty="0"/>'
           f'<a:t>‹#›</a:t></a:fld><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody></p:sp>')
    return ftr + num


def pic_xml(key, name, rId, x, y, w, h):
    pic = CT_Picture.new_pic(next_id(key), name, name, rId, emu(x), emu(y), emu(w), emu(h))
    return etree.tostring(pic).decode()


def gold_rule(key, x, y, w):
    """เส้นทองบางใต้หัวเรื่อง — ไล่เฉดทอง→โปร่ง ตามแบบ .ice-rule ของ Design System"""
    g = (f'<a:gradFill><a:gsLst><a:gs pos="0"><a:srgbClr val="{C["gold"]}"/></a:gs>'
         f'<a:gs pos="70000"><a:srgbClr val="{C["gold"]}"><a:alpha val="28000"/></a:srgbClr></a:gs>'
         f'<a:gs pos="100000"><a:srgbClr val="{C["gold"]}"><a:alpha val="0"/></a:srgbClr></a:gs></a:gsLst>'
         f'<a:lin ang="0" scaled="0"/></a:gradFill>')
    return cxn_line(key, "Gold Rule", x, y, x + w, y, C["gold"], 1.5, grad=g)


def hairline(key, x1, y1, x2, y2, alpha=22):
    return cxn_line(key, "Gold Hairline", x1, y1, x2, y2, C["gold"], 0.5, alpha=alpha)


# ─────────────────────────────────────────────────────────────────────────────
# ตำแหน่งกริด (นิ้ว) — safe margin 0.6 ทุกด้าน (กติกา ≥0.4)
# ─────────────────────────────────────────────────────────────────────────────
ML, MR = 0.6, 0.6
CW = W_IN - ML - MR              # 12.133
TITLE_Y, TITLE_H = 0.45, 1.05
ICON_X, ICON_Y, ICON_S = ML, 0.6, 0.75
TITLE_X = 1.55
LOGO_H = 0.45
LOGO_W = LOGO_H * LOGO_RATIO
LOGO_X = W_IN - MR - LOGO_W
TITLE_W = LOGO_X - 0.25 - TITLE_X
RULE_Y = 1.65
BODY_Y, BODY_H = 1.95, 4.55
FOOT_LINE_Y = 6.72


def title_block(key, prompt="หัวเรื่องเป็นประโยคสรุปหนึ่งข้อความต่อหน้า"):
    """หัวเรื่อง (title) + icon นำหน้า (pic idx 10) — ชุดเดียวกันทุกหน้าเนื้อหา"""
    return (ph_pic(key, "Title Icon Placeholder", 10, ICON_X, ICON_Y, ICON_S, ICON_S, prompt="icon")
            + ph(key, "Title Placeholder", "title", None, TITLE_X, TITLE_Y, TITLE_W, TITLE_H, 26, TITLE_INK,
                 bold=True, anchor="ctr", prompt=prompt))


# ─────────────────────────────────────────────────────────────────────────────
# MASTER (หน้าเนื้อหา = พื้นอ่อน)
# ─────────────────────────────────────────────────────────────────────────────
def master_xml(logo_rid):
    key = "master"
    g = C["grad_light"]
    bg = (f'<p:bg><p:bgPr>{grad_fill([(0, g[0], None), (55, g[1], None), (100, g[2], None)])}'
          f'<a:effectLst/></p:bgPr></p:bg>')
    shapes = ''
    # โลโก้สี มุมขวาบน (กติกา CI: ใช้โลโก้สีบนพื้นอ่อน วางมุมขวาบน มีที่ว่างรอบ)
    shapes += pic_xml(key, "iCE Logo", logo_rid, LOGO_X, 0.38, LOGO_W, LOGO_H)
    # เส้นทองบางเหนือ footer (ตกแต่งอย่างเดียว — ทองไม่ใช้กับข้อความ)
    shapes += hairline(key, ML, FOOT_LINE_Y, W_IN - MR, FOOT_LINE_Y, alpha=30)
    # placeholder ระดับ master (layout สืบทอดรูปแบบจากตรงนี้)
    shapes += ph(key, "Title Placeholder", "title", None, TITLE_X, TITLE_Y, TITLE_W, TITLE_H, 26, TITLE_INK,
                 bold=True, anchor="ctr", prompt="หัวเรื่องเป็นประโยคสรุป")
    shapes += ph(key, "Body Placeholder", "body", 1, ML, BODY_Y, CW, BODY_H, 18, BODY_INK, bullets=True,
                 prompt="เนื้อหา")
    shapes += footer_phs(key)

    tx = (f'<p:txStyles>'
          f'<p:titleStyle><a:lvl1pPr algn="l"><a:defRPr sz="2600" b="1">{solid(TITLE_INK)}{rpr_font()}</a:defRPr></a:lvl1pPr></p:titleStyle>'
          f'<p:bodyStyle>'
          f'<a:lvl1pPr marL="228600" indent="-228600"><a:spcBef><a:spcPts val="600"/></a:spcBef>'
          f'<a:buClr><a:srgbClr val="{C["teal"]}"/></a:buClr><a:buFont typeface="{FONT}"/><a:buChar char="•"/>'
          f'<a:defRPr sz="1800">{solid(BODY_INK)}{rpr_font()}</a:defRPr></a:lvl1pPr>'
          f'<a:lvl2pPr marL="457200" indent="-228600"><a:spcBef><a:spcPts val="300"/></a:spcBef>'
          f'<a:buClr><a:srgbClr val="{C["tealBright"]}"/></a:buClr><a:buFont typeface="{FONT}"/><a:buChar char="–"/>'
          f'<a:defRPr sz="1600">{solid(BODY_INK)}{rpr_font()}</a:defRPr></a:lvl2pPr>'
          f'<a:lvl3pPr marL="685800" indent="-228600"><a:buClr><a:srgbClr val="{C["tealBright"]}"/></a:buClr>'
          f'<a:buFont typeface="{FONT}"/><a:buChar char="–"/>'
          f'<a:defRPr sz="1600">{solid(BODY_INK)}{rpr_font()}</a:defRPr></a:lvl3pPr>'
          f'</p:bodyStyle>'
          f'<p:otherStyle><a:defPPr><a:defRPr lang="th-TH"/></a:defPPr>'
          f'<a:lvl1pPr algn="l"><a:defRPr sz="1400">{solid(BODY_INK)}{rpr_font()}</a:defRPr></a:lvl1pPr></p:otherStyle>'
          f'</p:txStyles>')
    return f'<p:cSld name="iCE-Propose">{bg}<p:spTree>{grp_hdr()}{shapes}</p:spTree></p:cSld>', tx


def grp_hdr():
    return ('<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
            '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>')


# ─────────────────────────────────────────────────────────────────────────────
# LAYOUTS
# ─────────────────────────────────────────────────────────────────────────────
def dark_bg():
    g = C["grad_dark"]
    return (f'<p:bg><p:bgPr>{grad_fill([(0, g[0], None), (45, g[1], None), (100, g[2], None)])}'
            f'<a:effectLst/></p:bgPr></p:bg>')


def dark_decor(key):
    """เส้นทองบางเฉียงสองเส้น (ลายเส้นทองจาง ๆ ตามหลักพื้นผิวแบรนด์ของ Design System)"""
    s = hairline(key, 8.4, 0.0, W_IN, 3.9, alpha=20)
    s += hairline(key, 9.6, 0.0, W_IN, 2.6, alpha=12)
    s += hairline(key, 0.0, 5.3, 4.6, H_IN, alpha=14)
    return s


def layout_cover(key, logo_w_rid):
    s = dark_decor(key)
    s += ph(key, "Kicker Placeholder", "body", 1, ML, 1.95, 9.0, 0.45, 14, C["goldLight"], anchor="b",
            prompt="ชื่อชุดเอกสาร เช่น ข้อเสนอโครงการ", autofit=False)
    s += ph(key, "Title Placeholder", "ctrTitle", None, ML, 2.5, 10.2, 1.7, 40, C["white"], bold=True,
            anchor="ctr", prompt="ชื่อเรื่องหลัก")
    s += ph(key, "Subtitle Placeholder", "subTitle", 2, ML, 4.3, 10.2, 0.85, 20, "D9E4EE", anchor="t",
            prompt="คำอธิบายสั้นหนึ่งประโยค")
    s += gold_rule(key, ML, 5.35, 2.2)
    s += ph(key, "Meta Placeholder", "body", 3, ML, 5.5, 9.0, 0.85, 14, "D9E4EE", anchor="t",
            prompt="จัดทำสำหรับ … · จัดทำโดย iCE Consulting · วันที่ …", autofit=False)
    s += pic_xml(key, "iCE Logo White", logo_w_rid, ML, H_IN - 0.55 - 0.42, 0.42 * LOGO_RATIO, 0.42)
    return dark_bg(), s, True


def layout_divider(key, logo_w_rid):
    s = dark_decor(key)
    s += pic_xml(key, "iCE Logo White", logo_w_rid, LOGO_X, 0.38, LOGO_W, LOGO_H)
    s += ph(key, "Section Number Placeholder", "body", 1, ML, 2.25, 2.3, 1.5, 60, C["goldLight"], anchor="ctr",
            align="l", prompt="01", autofit=False)
    s += ph(key, "Title Placeholder", "title", None, 3.05, 2.35, 9.7, 1.3, 36, C["white"], bold=True,
            anchor="ctr", prompt="ชื่อหัวข้อของช่วงนี้")
    s += ph(key, "Subtitle Placeholder", "body", 2, 3.05, 3.7, 9.7, 0.9, 20, "D9E4EE", anchor="t",
            prompt="หนึ่งประโยคบอกว่าช่วงนี้ตอบคำถามอะไร", autofit=False)
    s += gold_rule(key, 3.05, 4.8, 2.2)
    s += footer_phs(key, dark=True)
    return dark_bg(), s, True


def layout_action_title_body(key):
    s = title_block(key)
    s += gold_rule(key, ML, RULE_Y, 2.2)
    s += ph(key, "Body Placeholder", "body", 1, ML, BODY_Y, 7.0, BODY_H, 18, BODY_INK, bullets=True,
            prompt="ประเด็นสำคัญเป็น bullet ทีละข้อ")
    s += ph_pic(key, "Picture Placeholder", 2, 8.0, BODY_Y, W_IN - MR - 8.0, BODY_H, prompt="ภาพ / แผนภาพ / icon")
    s += footer_phs(key)
    return None, s, False


def layout_two_column(key):
    s = title_block(key)
    s += gold_rule(key, ML, RULE_Y, 2.2)
    colw = (CW - 0.55) / 2                      # 5.79
    lx, rx = ML, ML + colw + 0.55
    s += ph_pic(key, "Left Icon Placeholder", 5, lx, BODY_Y, 0.5, 0.5, prompt="icon")
    s += ph(key, "Left Heading Placeholder", "body", 1, lx + 0.65, BODY_Y - 0.05, colw - 0.65, 0.6, 20,
            C["navy"], bold=True, anchor="ctr", prompt="หัวข้อคอลัมน์ซ้าย", autofit=False)
    s += ph(key, "Left Body Placeholder", "body", 2, lx, BODY_Y + 0.7, colw, BODY_H - 0.7, 16, BODY_INK,
            bullets=True, prompt="เนื้อหาคอลัมน์ซ้าย")
    s += cxn_line(key, "Column Divider", lx + colw + 0.275, BODY_Y, lx + colw + 0.275, BODY_Y + BODY_H,
                  C["gold"], 0.75, alpha=35)
    s += ph_pic(key, "Right Icon Placeholder", 6, rx, BODY_Y, 0.5, 0.5, prompt="icon")
    s += ph(key, "Right Heading Placeholder", "body", 3, rx + 0.65, BODY_Y - 0.05, colw - 0.65, 0.6, 20,
            C["navy"], bold=True, anchor="ctr", prompt="หัวข้อคอลัมน์ขวา", autofit=False)
    s += ph(key, "Right Body Placeholder", "body", 4, rx, BODY_Y + 0.7, colw, BODY_H - 0.7, 16, BODY_INK,
            bullets=True, prompt="เนื้อหาคอลัมน์ขวา")
    s += footer_phs(key)
    return None, s, False


def layout_three_card(key):
    s = title_block(key)
    s += gold_rule(key, ML, RULE_Y, 2.2)
    gap = 0.3
    cw = (CW - 2 * gap) / 3                     # 3.844
    strip = [C["navy"], C["tealBright"], C["teal"]]   # ไล่เฉดน้ำเงิน→ฟ้า ข้ามการ์ดสามใบ
    for i in range(3):
        x = ML + i * (cw + gap)
        s += sp_rect(key, f"Card {i + 1}", x, BODY_Y, cw, BODY_H, solid(CARD_BG), ln(CARD_LINE, 0.75),
                     prst="roundRect")
        s += sp_rect(key, f"Card Strip {i + 1}", x + 0.02, BODY_Y, cw - 0.04, 0.09,
                     grad_fill([(0, strip[i], None), (100, strip[min(i + 1, 2)], None)], 90), no_ln())
        base = i * 3
        s += ph_pic(key, f"Card {i + 1} Icon Placeholder", base + 1, x + 0.3, BODY_Y + 0.35, 0.7, 0.7, prompt="icon")
        s += ph(key, f"Card {i + 1} Heading Placeholder", "body", base + 2, x + 0.3, BODY_Y + 1.2, cw - 0.6, 0.6, 18,
                C["navy"], bold=True, anchor="ctr", prompt=f"หัวข้อการ์ดที่ {i + 1}", autofit=False)
        s += ph(key, f"Card {i + 1} Body Placeholder", "body", base + 3, x + 0.3, BODY_Y + 1.9, cw - 0.6,
                BODY_H - 2.2, 16, BODY_INK, bullets=True, prompt="เนื้อหาของการ์ด")
    s += footer_phs(key)
    return None, s, False


def layout_table(key):
    s = title_block(key)
    s += gold_rule(key, ML, RULE_Y, 2.2)
    s += ph_tbl(key, "Table Placeholder", 1, ML, BODY_Y, CW, 3.9)
    s += ph(key, "Note Placeholder", "body", 2, ML, 6.0, CW, 0.5, 12, MUTED, anchor="ctr",
            prompt="หมายเหตุหรือที่มาของข้อมูล", autofit=False)
    s += footer_phs(key)
    return None, s, False


def layout_timeline(key):
    s = title_block(key)
    s += gold_rule(key, ML, RULE_Y, 2.2)
    n = 4
    seg = CW / n
    track_y = 3.7
    g = C["grad_brand"]
    s += cxn_line(key, "Timeline Track", ML + 0.3, track_y, W_IN - MR - 0.3, track_y, C["navy"], 5,
                  grad=grad_fill([(0, g[0], None), (100, g[1], None)], 90))
    tones = [C["navyDeep"], C["navy"], C["tealBright"], C["teal"]]
    for i in range(n):
        cx = ML + (i + 0.5) * seg
        d = 0.52
        s += sp_rect(key, f"Milestone {i + 1}", cx - d / 2, track_y - d / 2, d, d, solid(tones[i]),
                     ln(C["white"], 1.5), prst="ellipse",
                     text_xml=para(str(i + 1), 14, C["white"], bold=True, align="ctr"))
        s += ph(key, f"Phase {i + 1} Label Placeholder", "body", i + 1, cx - seg / 2 + 0.1, 2.55, seg - 0.2, 0.8, 16,
                C["navy"], bold=True, align="ctr", anchor="b", prompt=f"ช่วงที่ {i + 1} · ระยะเวลา", autofit=False)
        s += ph(key, f"Phase {i + 1} Detail Placeholder", "body", n + i + 1, cx - seg / 2 + 0.1, 4.1, seg - 0.2, 2.4, 16,
                BODY_INK, align="ctr", anchor="t", prompt="สิ่งที่เกิดขึ้นในช่วงนี้")
    s += footer_phs(key)
    return None, s, False


def layout_closing(key, logo_w_rid):
    s = dark_decor(key)
    s += ph(key, "Title Placeholder", "ctrTitle", None, ML, 2.4, 10.2, 1.4, 40, C["white"], bold=True,
            anchor="ctr", prompt="ขอบคุณ")
    s += ph(key, "Subtitle Placeholder", "subTitle", 1, ML, 3.9, 10.2, 0.8, 20, "D9E4EE", anchor="t",
            prompt="ประโยคปิดหรือขั้นถัดไปที่เสนอ")
    s += gold_rule(key, ML, 4.95, 2.2)
    s += ph(key, "Contact Placeholder", "body", 2, ML, 5.1, 9.0, 1.0, 14, "D9E4EE", anchor="t",
            prompt="ชื่อผู้ติดต่อ · อีเมล · โทรศัพท์", autofit=False)
    s += pic_xml(key, "iCE Logo White", logo_w_rid, ML, H_IN - 0.55 - 0.42, 0.42 * LOGO_RATIO, 0.42)
    return dark_bg(), s, True


def layout_appendix(key):
    s = sp_rect(key, "Appendix Tag", TITLE_X, 0.22, 1.4, 0.26, solid(C["gold"], 18), no_ln(), prst="roundRect",
                text_xml=para("ภาคผนวก", 10, C["goldDeep"], bold=True, align="ctr"))
    s += title_block(key, prompt="หัวเรื่องของภาคผนวก")
    s += gold_rule(key, ML, RULE_Y, 2.2)
    s += ph(key, "Body Placeholder", "body", 1, ML, BODY_Y, CW, BODY_H, 16, BODY_INK, bullets=True,
            prompt="รายละเอียดประกอบ")
    s += footer_phs(key)
    return None, s, False


def layout_blank(key):
    """หน้าว่างบนแม่แบบ — มีแค่ของ master (พื้นอ่อน โลโก้ เส้นทอง footer เลขหน้า)
    ใช้กับหน้าที่ต้องวาด infographic เองทั้งหน้า และเป็นหน้าที่ build_deck.py ใช้ในโหมดแม่แบบ"""
    return None, footer_phs(key), False


LAYOUT_ORDER = ["cover", "divider", "action-title-body", "two-column", "three-card",
                "table", "timeline", "closing", "appendix", "blank"]


# ─────────────────────────────────────────────────────────────────────────────
# THEME
# ─────────────────────────────────────────────────────────────────────────────
def patch_theme(theme_part):
    root = etree.fromstring(theme_part.blob)
    a = NS["a"]
    clr = root.find(f".//{{{a}}}clrScheme")
    scheme = {
        "dk1": C["ink"], "lt1": C["white"], "dk2": C["navyDeep"], "lt2": C["grad_light"][1],
        "accent1": C["navy"], "accent2": C["teal"], "accent3": C["tealBright"], "accent4": C["gold"],
        "accent5": C["success"], "accent6": C["warning"], "hlink": C["navy"], "folHlink": C["teal"],
    }
    clr.set("name", "iCE")
    for child in list(clr):
        tag = etree.QName(child).localname
        if tag in scheme:
            for gc in list(child):
                child.remove(gc)
            el = etree.SubElement(child, f"{{{a}}}srgbClr")
            el.set("val", scheme[tag])
    font_scheme = root.find(f".//{{{a}}}fontScheme")
    font_scheme.set("name", "iCE")
    for fs in ("majorFont", "minorFont"):
        node = font_scheme.find(f"{{{a}}}{fs}")
        for child in list(node):
            node.remove(child)
        for slot in ("latin", "ea", "cs"):
            el = etree.SubElement(node, f"{{{a}}}{slot}")
            el.set("typeface", FONT)
        # ⭐ ผูกสคริปต์ไทย/ญี่ปุ่นฯ ตรง ๆ ด้วย — บางโปรแกรมอ่านตาราง a:font ก่อน ea/cs
        for script in ("Thai", "Jpan", "Hans", "Hant", "Arab", "Hebr"):
            el = etree.SubElement(node, f"{{{a}}}font")
            el.set("script", script)
            el.set("typeface", FONT)
    theme_part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


# ─────────────────────────────────────────────────────────────────────────────
# ประกอบไฟล์
# ─────────────────────────────────────────────────────────────────────────────
def replace_cSld(element, cSld_xml):
    p = NS["p"]
    old = element.find(f"{{{p}}}cSld")
    new = parse_xml(f'<p:cSld {NSDECL}>{cSld_xml}</p:cSld>') if not cSld_xml.startswith("<p:cSld") \
        else parse_xml(cSld_xml.replace("<p:cSld", f"<p:cSld {NSDECL}", 1))
    element.replace(old, new)


def build(out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    master = prs.slide_master

    # ── theme ──
    theme_part = master.part.part_related_by(RT.THEME)
    patch_theme(theme_part)

    # ── master ──
    _, logo_rid = master.part.get_or_add_image_part(LOGO_COLOR)
    cSld, tx = master_xml(logo_rid)
    replace_cSld(master._element, cSld)
    p = NS["p"]
    old_tx = master._element.find(f"{{{p}}}txStyles")
    master._element.replace(old_tx, parse_xml(tx.replace("<p:txStyles>", f"<p:txStyles {NSDECL}>", 1)))

    # ── layouts: เขียนทับ 9 ตัวแรก แล้วลบที่เหลือ ──
    layouts = list(prs.slide_layouts)
    for lay in layouts[len(LAYOUT_ORDER):]:
        prs.slide_layouts.remove(lay)
    for name, lay in zip(LAYOUT_ORDER, layouts):
        key = f"layout:{name}"
        _ids[key] = 1
        # ล้าง relationship ไปยังรูปเก่า (layout ค่าเริ่มต้นไม่มี แต่กันไว้เมื่อรันซ้ำ)
        need_white = name in ("cover", "divider", "closing")
        logo_w_rid = lay.part.get_or_add_image_part(LOGO_WHITE)[1] if need_white else None
        fn = {
            "cover": lambda: layout_cover(key, logo_w_rid),
            "divider": lambda: layout_divider(key, logo_w_rid),
            "action-title-body": lambda: layout_action_title_body(key),
            "two-column": lambda: layout_two_column(key),
            "three-card": lambda: layout_three_card(key),
            "table": lambda: layout_table(key),
            "timeline": lambda: layout_timeline(key),
            "closing": lambda: layout_closing(key, logo_w_rid),
            "appendix": lambda: layout_appendix(key),
            "blank": lambda: layout_blank(key),
        }[name]
        bg, shapes, dark = fn()
        cSld = f'<p:cSld name="{name}">{bg or ""}<p:spTree>{grp_hdr()}{shapes}</p:spTree></p:cSld>'
        replace_cSld(lay._element, cSld)
        el = lay._element
        # หน้าพื้นเข้มซ่อนรูปทรงของ master (โลโก้สี/เส้นทองของพื้นอ่อน) แล้วใส่ของตัวเอง
        if dark:
            el.set("showMasterSp", "0")
        elif "showMasterSp" in el.attrib:
            del el.attrib["showMasterSp"]
        el.set("preserve", "1")
        if "type" in el.attrib:
            del el.attrib["type"]
        # matchingName ช่วยให้ PowerPoint จับคู่ layout ตอน copy ข้ามไฟล์
        el.set("matchingName", name)

    # เอกสารประกอบของไฟล์
    prs.core_properties.title = "iCE-Propose Master"
    prs.core_properties.author = "iCE Consulting"
    prs.core_properties.subject = "แม่แบบ .pptx ของ iCE — 9 layout · 16:9"
    prs.core_properties.comments = f"สร้างโดย make_master.py · ฟอนต์ theme = {FONT} · สีจาก iCE Design System tokens.json"

    prs.save(out_path)
    names = [l.name for l in prs.slide_layouts]
    print(f"OK: {out_path}")
    print(f"   layouts ({len(names)}): {', '.join(names)}")
    print(f"   theme font (latin/ea/cs): {FONT}")
    print(f"   slide size: {prs.slide_width / EMU:.3f} × {prs.slide_height / EMU:.2f} นิ้ว")


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "iCE-Propose_Master.pptx")
    build(out)
