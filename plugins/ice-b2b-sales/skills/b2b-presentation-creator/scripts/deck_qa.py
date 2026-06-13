#!/usr/bin/env python3

"""
deck_qa.py
Automated QA on generated .pptx — converts to images, runs grep checks,
generates markdown report with Critical/Moderate issue detection.
"""

import argparse
import logging
import re
import shutil
import subprocess
import sys
from pathlib import Path
from typing import List, Tuple

logging.basicConfig(level=logging.INFO, format="[%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

CRITICAL_PATTERNS = [
    (r"xxxx", "Placeholder (xxxx)"),
    (r"lorem\s+ipsum", "Lorem ipsum placeholder"),
    (r"\[customer\]", "Unreplaced [customer] tag"),
    (r"\[date\]", "Unreplaced [date] tag"),
    (r"\[name\]", "Unreplaced [name] tag"),
    (r"as\s+an?\s+(AI|advisor|expert)", "AI scaffolding phrase"),
    (r"the\s+following\s+bullet", "AI scaffolding: 'the following bullet'"),
    (r"in\s+this\s+slide", "AI scaffolding: 'in this slide'"),
]

MODERATE_PATTERNS = [
    (r"tbd", "To-be-determined (TBD)"),
    (r"placeholder", "Generic placeholder text"),
    (r"expand\s+on", "Incomplete instruction"),
]

# PPTX Lesson #18: chars that make PowerPoint for Mac REJECT the whole file (Repair)
# but LibreOffice/qlmanage pass (false-green). Detected via text-extract (NO LibreOffice
# dependency) so the check works even when no render engine is installed.
CHAR_BLOCKLIST = {
    "→": "▸",  # U+2192 RIGHTWARDS ARROW         → ▸ (U+25B8, PowerPoint-safe)
    "⟶": "▸",  # U+27F6 LONG RIGHTWARDS ARROW    → ▸
    "➜": "▸",  # U+2799 HEAVY ROUND-TIPPED ARROW → ▸
    "➔": "▸",  # U+2794 HEAVY WIDE-HEADED ARROW  → ▸
    "➙": "▸",  # U+2799 HEAVY RIGHTWARDS ARROW   → ▸
}


def check_libreoffice() -> bool:
    """Check if LibreOffice is available."""
    return shutil.which("soffice") is not None


def pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    """Convert PPTX to PDF using LibreOffice (preview render only — NOT a validation pass).
    Returns True if rendered, False if LibreOffice absent (char/content checks still run).
    NOTE: LibreOffice is false-green — it cannot see corruption/16:9/General-Failure/U+2192
    that real PowerPoint rejects. Final visual validation must open in REAL PowerPoint."""
    if not check_libreoffice():
        logger.warning(
            "LibreOffice not found — skipping preview render (this is OK). "
            "Text/char checks still run. Final visual check: open in REAL PowerPoint "
            "(LibreOffice render = false-green, install only for preview: brew install libreoffice)."
        )
        return False

    try:
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(Path(pdf_path).parent),
                pptx_path,
            ],
            check=True,
            capture_output=True,
            timeout=120,
        )
        logger.info(f"PDF generated: {pdf_path}")
        return True
    except subprocess.CalledProcessError as e:
        logger.warning(f"LibreOffice preview render failed (non-fatal): {e.stderr.decode()}")
        return False


def pdf_to_images(pdf_path: str, output_dir: str, dpi: int) -> List[str]:
    """Convert PDF to JPG images using pdftoppm."""
    try:
        subprocess.run(
            [
                "pdftoppm",
                "-jpeg",
                "-singlefile",
                f"-r {dpi}",
                pdf_path,
                str(Path(output_dir) / "slide"),
            ],
            check=True,
            capture_output=True,
        )
    except FileNotFoundError:
        logger.error(
            "pdftoppm not found. Install: brew install poppler (macOS) "
            "or apt-get install poppler-utils (Linux)"
        )
        sys.exit(1)
    except subprocess.CalledProcessError as e:
        logger.error(f"pdftoppm conversion failed: {e.stderr.decode()}")
        sys.exit(1)

    images = sorted(Path(output_dir).glob("slide-*.jpg"))
    logger.info(f"Generated {len(images)} slide images")
    return [str(img) for img in images]


def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed. Install: pip install python-pptx")
        sys.exit(1)

    try:
        prs = Presentation(pptx_path)
        text_blocks = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_blocks.append(shape.text)
        return "\n".join(text_blocks)
    except Exception as e:
        logger.error(f"Failed to extract text from PPTX: {e}")
        sys.exit(1)


def find_issues(text: str, patterns: List[Tuple[str, str]]) -> List[str]:
    """Find issues in text using regex patterns."""
    issues = []
    for pattern, description in patterns:
        if re.search(pattern, text, re.IGNORECASE):
            issues.append(description)
    return issues


def find_forbidden_chars(pptx_path: str) -> List[str]:
    """Lesson #18: scan each slide for PowerPoint-rejecting chars (→ etc).
    Returns per-slide issue strings. Uses python-pptx — no LibreOffice needed."""
    try:
        from pptx import Presentation
    except ImportError:
        return []
    issues = []
    try:
        prs = Presentation(pptx_path)
        for idx, slide in enumerate(prs.slides, 1):
            stext = "\n".join(sh.text for sh in slide.shapes if hasattr(sh, "text"))
            for bad, good in CHAR_BLOCKLIST.items():
                if bad in stext:
                    issues.append(
                        f"Slide {idx}: forbidden char '{bad}' (U+{ord(bad):04X}) "
                        f"→ PowerPoint rejects whole file (Repair). Replace with '{good}'."
                    )
    except Exception as e:
        logger.error(f"forbidden-char scan failed: {e}")
    return issues


def generate_report(
    pptx_path: str,
    output_dir: str,
    critical_issues: List[str],
    moderate_issues: List[str],
    slide_count: int,
    images: List[str],
) -> str:
    """Generate markdown QA report."""
    report_lines = [
        "# Deck QA Report",
        "",
        f"**File:** {Path(pptx_path).name}",
        f"**Generated:** {Path(pptx_path).stat().st_mtime}",
        f"**Slide Count:** {slide_count}",
        "",
        "## Summary",
        f"- Critical Issues: {len(critical_issues)}",
        f"- Moderate Issues: {len(moderate_issues)}",
        "",
    ]

    if critical_issues:
        report_lines.extend(["## Critical Issues", ""])
        for issue in critical_issues:
            report_lines.append(f"- {issue}")
        report_lines.append("")

    if moderate_issues:
        report_lines.extend(["## Moderate Issues", ""])
        for issue in moderate_issues:
            report_lines.append(f"- {issue}")
        report_lines.append("")

    report_lines.extend(
        [
            "## Slide Images",
            f"Images saved to: `{output_dir}`",
            "",
            "## Next Steps",
            "- Review slide images for visual QA",
            "- Verify formatting and alignment",
            "- Consider manual review via visual QA subagent",
            "",
        ]
    )

    return "\n".join(report_lines)


def main() -> None:
    """Main entry point."""
    parser = argparse.ArgumentParser(
        description="Automated QA on PPTX presentations."
    )
    parser.add_argument("--deck", required=True, help="Path to PPTX file")
    parser.add_argument(
        "--output-dir", required=True, help="Directory for images and report"
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Fail on Moderate issues (default: Critical only)",
    )
    parser.add_argument("--dpi", type=int, default=150, help="DPI for image conversion")

    args = parser.parse_args()

    output_path = Path(args.output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    pptx_path = Path(args.deck)
    if not pptx_path.exists():
        logger.error(f"PPTX file not found: {args.deck}")
        sys.exit(1)

    logger.info(f"Starting QA on: {args.deck}")

    # Lesson #18: forbidden-char scan FIRST — runs always, no LibreOffice needed
    forbidden = find_forbidden_chars(str(pptx_path))

    text = extract_text_from_pptx(str(pptx_path))
    critical_issues = find_issues(text, CRITICAL_PATTERNS) + forbidden
    moderate_issues = find_issues(text, MODERATE_PATTERNS)

    # LibreOffice preview render (OPTIONAL — preview only, NOT a validation pass).
    # If absent, char/content checks above still cover correctness; visual check = real PowerPoint.
    images = []
    with_pdf = output_path / f"{pptx_path.stem}.pdf"
    if pptx_to_pdf(str(pptx_path), str(with_pdf)):
        images = pdf_to_images(str(with_pdf), str(output_path), args.dpi)
    else:
        logger.warning("No preview images — open the deck in REAL PowerPoint for visual QA.")

    try:
        from pptx import Presentation
        slide_count = len(Presentation(str(pptx_path)).slides)
    except Exception:
        slide_count = len(images)

    report = generate_report(
        str(pptx_path),
        str(output_path),
        critical_issues,
        moderate_issues,
        slide_count,
        images,
    )

    report_path = output_path / f"qa_{pptx_path.stem}.md"
    report_path.write_text(report)
    logger.info(f"QA report saved: {report_path}")

    print(report)

    exit_code = 0
    if critical_issues:
        logger.warning(f"Found {len(critical_issues)} CRITICAL issues")
        exit_code = 1
    if args.strict and moderate_issues:
        logger.warning(f"Found {len(moderate_issues)} MODERATE issues (strict mode)")
        exit_code = 2

    sys.exit(exit_code)


if __name__ == "__main__":
    main()
