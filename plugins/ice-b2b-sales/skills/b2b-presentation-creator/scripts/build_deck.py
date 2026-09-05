#!/usr/bin/env python3
"""
B2B Presentation Creator - PPTX Deck Builder
Builds a .pptx deck from an outline JSON, applying a theme and language mode.

Usage:
    python build_deck.py --outline <path> --theme <path> --language <th|en|bilingual> \
        --output <path> [--debug] [--validate-only]
"""

import argparse
import json
import logging
import os
import sys
from pathlib import Path
from typing import Dict, List, Tuple, Any, Optional

# ⭐ FONT POLICY มาจาก SSOT เดียว (ice-doc-builder §3.0) — ห้าม hard-code ชื่อฟอนต์ในไฟล์นี้
sys.path.insert(0, os.path.expanduser("~/.claude/agents/_lib"))
from font_policy import RAILS  # noqa: E402

RAIL = os.environ.get("ICE_RAIL", "private")
RAIL_FONT = RAILS[RAIL if RAIL in RAILS else "private"]["font"]
# คู่ละตินสำหรับ deck ภาษาอังกฤษล้วน — ไม่มีอักขระไทยจึงไม่อยู่ใต้ราง (§3.0 ครอบเฉพาะงานที่มีไทย)
LATIN_HEADER = "Inter"
LATIN_BODY = "Lora"
# ฟอนต์สำหรับ glyph สัญลักษณ์ล้วน (▸ ฯลฯ) — ไม่ใช่ข้อความ ไม่อยู่ใต้ราง
SYMBOL_FONT = "Arial"

# ⭐ V02 (2026.09.05 · Wave B) โหมดแม่แบบ: ICE_TEMPLATE=<path .pptx แม่แบบ iCE>
#   ไม่ตั้ง = พฤติกรรมเดิมทุกอย่าง 100% · ตั้งแล้ว = เปิดแม่แบบ (ฟอนต์ theme/สี/โลโก้/footer/เลขหน้า
#   มาจาก master) ลบสไลด์ตัวอย่างที่ติดมา แล้วเลือก layout ตามชื่อ:
#     title-hero → "cover" (พื้นเข้มของแม่แบบ ไม่วาดพื้นหลังเอง) · หน้าอื่น → "blank" ของแม่แบบ
#     หรือระบุเองรายหน้าใน outline ด้วย "template_layout": "<ชื่อ layout>"
#   ฟังก์ชันวาดเนื้อหาของไฟล์นี้ยังวาด textbox/รูปทรงเองเหมือนเดิม (ยังไม่เติม placeholder) —
#   placeholder ของ layout จึงถูกถอดออกก่อน ยกเว้น footer/เลขหน้า ซึ่งยกมาจาก layout แทน add_footer()
#   การเติมเนื้อหาลง placeholder จริงทำที่ ~/.claude/agents/_lib/build_pptx.py (โหมดแม่แบบ V02R05)
ICE_TEMPLATE = os.environ.get("ICE_TEMPLATE") or None
TEMPLATE_DEFAULT_LAYOUT = {"title-hero": "cover"}
TEMPLATE_FALLBACK_LAYOUT = "blank"


def _template_layout(prs, name: str):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    raise ValueError(f"แม่แบบไม่มี layout ชื่อ '{name}' (มี: {', '.join(l.name for l in prs.slide_layouts)})")


def _template_remove_slides(prs) -> None:
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        prs.part.drop_rel(sldId.rId)
        sldIdLst.remove(sldId)


def _template_prepare_slide(slide, layout) -> None:
    """ถอด placeholder ที่ clone มาจาก layout (เนื้อหาถูกวาดเองอยู่แล้ว) แล้วยก footer/เลขหน้าจาก layout"""
    import copy as _copy
    for shp in list(slide.placeholders):
        shp._element.getparent().remove(shp._element)
    for shp in layout.placeholders:
        if shp.placeholder_format.idx in (11, 12):
            slide.shapes._spTree.append(_copy.deepcopy(shp._element))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)]
)
logger = logging.getLogger(__name__)


# Constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.5)
TITLE_FONTSIZE = 44
SUBTITLE_FONTSIZE = 28
BODY_FONTSIZE = 16
SMALL_FONTSIZE = 14
STAT_FONTSIZE = 96


def hex_to_rgb(hex_str: str) -> RGBColor:
    """
    Convert hex color string to RGBColor.

    Args:
        hex_str: Color in format "#RRGGBB"

    Returns:
        RGBColor object

    Raises:
        ValueError: If hex string is invalid
    """
    hex_str = hex_str.lstrip("#")
    if len(hex_str) != 6:
        raise ValueError(f"Invalid hex color: {hex_str}")
    try:
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
    except ValueError as e:
        raise ValueError(f"Invalid hex color: {hex_str}") from e


def pick_font(theme: Dict[str, Any], role: str, language: str) -> str:
    """
    Select appropriate font based on language mode and role (header/body).

    Args:
        theme: Theme dictionary with fonts section
        role: Either 'header' or 'body'
        language: Either 'th', 'en', or 'bilingual'

    Returns:
        Font name string
    """
    fonts = theme.get("fonts", {})

    # ⭐ V02R01 (2026.08.04): default มาจาก SSOT font_policy.RAILS — ไม่ hard-code อีก
    #   เดิม default = "Sarabun" / "IBM Plex Sans Thai" (ไม่มี Looped = คนละ family, cut ไม่มีหัว)
    #   theme.json ยัง override ได้เหมือนเดิม = เคสลูกค้าบังคับ Latin brand font (§3.0 ②)
    if language == "th":
        return fonts.get(f"{role}TH", RAIL_FONT)
    elif language == "en":
        return fonts.get(f"{role}EN", LATIN_BODY if role == "body" else LATIN_HEADER)
    else:  # bilingual — ไทยครอบทั้งคู่ได้เพราะฟอนต์ราง single-family (มีละตินในตัว §3.0 ①)
        return fonts.get(f"{role}EN" if role == "header" else "bodyTH", RAIL_FONT)


def add_text_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    fontsize: int,
    font_name: str,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """
    Add a text box to a slide with specified properties.

    Args:
        slide: Slide object to add to
        left, top, width, height: Position and size in inches
        text: Text content
        fontsize: Font size in points
        font_name: Font family name
        color: RGBColor object
        bold: Whether text is bold
        align: Text alignment
    """
    if not text:
        return

    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(fontsize)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align


def add_rectangle_shape(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: RGBColor,
    line_color: Optional[RGBColor] = None,
) -> None:
    """
    Add a rectangle shape to a slide.

    Args:
        slide: Slide object
        left, top, width, height: Position and size in inches
        fill_color: Fill color as RGBColor
        line_color: Outline color (None for no outline)
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.color.rgb = fill_color


def add_footer(slide, theme: Dict[str, Any], page_num: int, total: int) -> None:
    """
    Add a thin footer with page number to a slide.

    Args:
        slide: Slide object
        theme: Theme dictionary for color
        page_num: Current page number (1-indexed)
        total: Total slide count
    """
    text = f"{page_num} / {total}"
    add_text_box(
        slide,
        left=0.5,
        top=7.0,
        width=12.333,
        height=0.4,
        text=text,
        fontsize=11,
        font_name=LATIN_HEADER,
        color=hex_to_rgb(theme["colors"]["muted"]),
        align=PP_ALIGN.RIGHT,
    )


def add_title_hero(
    slide, theme: Dict[str, Any], content: Dict[str, Any], language: str
) -> None:
    """
    Add a title-hero slide (typically cover page with large title and subtitle).

    Args:
        slide: Slide object
        theme: Theme dictionary
        content: Slide content with 'title', 'subtitle', optionally 'title_th', 'subtitle_th'
        language: 'th', 'en', or 'bilingual'
    """
    colors = theme["colors"]
    bg_color = hex_to_rgb(colors["background"])
    primary_color = hex_to_rgb(colors["primary"])
    text_color = hex_to_rgb(colors["text"])

    if ICE_TEMPLATE:
        # โหมดแม่แบบ: layout "cover" มีพื้นเข้มไล่เฉด + โลโก้ขาวอยู่แล้ว → ไม่วาดพื้นหลังทับ · ตัวอักษรขาว
        primary_color = text_color = RGBColor(0xFF, 0xFF, 0xFF)
    else:
        # Background
        add_rectangle_shape(slide, 0, 0, 13.333, 7.5, bg_color)

        # Accent bar on the left
        add_rectangle_shape(slide, 0, 0, 0.1, 7.5, primary_color)

    # Title
    title = content.get("title") or content.get("title_en", "")
    if language == "th" and "title_th" in content:
        title = content["title_th"]

    add_text_box(
        slide,
        left=0.75,
        top=2.5,
        width=11.833,
        height=1.5,
        text=title,
        fontsize=TITLE_FONTSIZE,
        font_name=pick_font(theme, "header", language),
        color=primary_color,
        bold=True,
    )

    # Subtitle
    subtitle = content.get("subtitle") or content.get("subtitle_en", "")
    if language == "th" and "subtitle_th" in content:
        subtitle = content["subtitle_th"]

    if subtitle:
        add_text_box(
            slide,
            left=0.75,
            top=4.2,
            width=11.833,
            height=1.0,
            text=subtitle,
            fontsize=SUBTITLE_FONTSIZE,
            font_name=pick_font(theme, "body", language),
            color=text_color,
        )


def add_two_column(
    slide, theme: Dict[str, Any], content: Dict[str, Any], language: str
) -> None:
    """
    Add a two-column layout slide.

    Args:
        slide: Slide object
        theme: Theme dictionary
        content: Slide content with 'title', 'left' (list), 'right' (dict or list)
        language: Language mode
    """
    colors = theme["colors"]
    bg_color = hex_to_rgb(colors["background"])
    primary_color = hex_to_rgb(colors["primary"])
    text_color = hex_to_rgb(colors["text"])

    add_rectangle_shape(slide, 0, 0, 13.333, 7.5, bg_color)

    # Title
    title = content.get("title") or content.get("title_en", "")
    if language == "th" and "title_th" in content:
        title = content["title_th"]

    add_text_box(
        slide,
        left=0.5,
        top=0.5,
        width=12.333,
        height=0.6,
        text=title,
        fontsize=TITLE_FONTSIZE - 8,
        font_name=pick_font(theme, "header", language),
        color=primary_color,
        bold=True,
    )

    # Left column
    left_items = content.get("left", [])
    if isinstance(left_items, list):
        left_text = "\n".join([str(item) for item in left_items if item])
    else:
        left_text = str(left_items) if left_items else ""

    add_text_box(
        slide,
        left=0.75,
        top=1.3,
        width=5.833,
        height=5.7,
        text=left_text,
        fontsize=BODY_FONTSIZE,
        font_name=pick_font(theme, "body", language),
        color=text_color,
    )

    # Right column
    right_items = content.get("right", [])
    if isinstance(right_items, dict):
        right_text = json.dumps(right_items, indent=2, ensure_ascii=False)
    elif isinstance(right_items, list):
        right_text = "\n".join([str(item) for item in right_items if item])
    else:
        right_text = str(right_items) if right_items else ""

    add_text_box(
        slide,
        left=6.75,
        top=1.3,
        width=5.833,
        height=5.7,
        text=right_text,
        fontsize=BODY_FONTSIZE,
        font_name=pick_font(theme, "body", language),
        color=text_color,
    )


def add_stat_callout(
    slide, theme: Dict[str, Any], content: Dict[str, Any], language: str
) -> None:
    """
    Add a stat callout slide with large number, label, and context.

    Args:
        slide: Slide object
        theme: Theme dictionary
        content: Slide content with 'value', 'label', 'context'
        language: Language mode
    """
    colors = theme["colors"]
    bg_color = hex_to_rgb(colors["background"])
    accent_color = hex_to_rgb(colors["accent"])
    text_color = hex_to_rgb(colors["text"])
    primary_color = hex_to_rgb(colors["primary"])

    add_rectangle_shape(slide, 0, 0, 13.333, 7.5, bg_color)

    # Accent background for stat
    add_rectangle_shape(slide, 0, 0, 13.333, 4.0, hex_to_rgb(colors["secondary"]))

    # Large stat value
    value = content.get("value", "")
    add_text_box(
        slide,
        left=0.75,
        top=1.2,
        width=11.833,
        height=2.0,
        text=value,
        fontsize=STAT_FONTSIZE,
        font_name=pick_font(theme, "header", language),
        color=accent_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Label
    label = content.get("label") or content.get("label_en", "")
    if language == "th" and "label_th" in content:
        label = content["label_th"]

    add_text_box(
        slide,
        left=0.75,
        top=4.5,
        width=11.833,
        height=0.8,
        text=label,
        fontsize=SUBTITLE_FONTSIZE,
        font_name=pick_font(theme, "header", language),
        color=primary_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Context
    context = content.get("context") or content.get("context_en", "")
    if language == "th" and "context_th" in content:
        context = content["context_th"]

    if context:
        add_text_box(
            slide,
            left=0.75,
            top=5.5,
            width=11.833,
            height=1.5,
            text=context,
            fontsize=BODY_FONTSIZE,
            font_name=pick_font(theme, "body", language),
            color=text_color,
            align=PP_ALIGN.CENTER,
        )


def add_icon_text_rows(
    slide, theme: Dict[str, Any], content: Dict[str, Any], language: str
) -> None:
    """
    Add an icon-text-rows slide with icon, header, and body for each row.

    Args:
        slide: Slide object
        theme: Theme dictionary
        content: Slide content with 'title' and 'rows' (list of {icon, header, body})
        language: Language mode
    """
    colors = theme["colors"]
    bg_color = hex_to_rgb(colors["background"])
    primary_color = hex_to_rgb(colors["primary"])
    text_color = hex_to_rgb(colors["text"])
    accent_color = hex_to_rgb(colors["accent"])

    add_rectangle_shape(slide, 0, 0, 13.333, 7.5, bg_color)

    # Title
    title = content.get("title") or content.get("title_en", "")
    if language == "th" and "title_th" in content:
        title = content["title_th"]

    add_text_box(
        slide,
        left=0.5,
        top=0.5,
        width=12.333,
        height=0.6,
        text=title,
        fontsize=TITLE_FONTSIZE - 8,
        font_name=pick_font(theme, "header", language),
        color=primary_color,
        bold=True,
    )

    # Rows
    rows = content.get("rows", [])
    row_height = 1.4
    start_top = 1.3

    for idx, row in enumerate(rows[:4]):  # Limit to 4 rows per slide
        top = start_top + (idx * row_height)

        # Icon placeholder (small circle)
        add_rectangle_shape(
            slide,
            left=0.75,
            top=top,
            width=0.3,
            height=0.3,
            fill_color=accent_color,
        )

        # Header
        header = row.get("header") or row.get("header_en", "")
        if language == "th" and "header_th" in row:
            header = row["header_th"]

        add_text_box(
            slide,
            left=1.2,
            top=top,
            width=11.333,
            height=0.4,
            text=header,
            fontsize=BODY_FONTSIZE,
            font_name=pick_font(theme, "header", language),
            color=primary_color,
            bold=True,
        )

        # Body
        body = row.get("body") or row.get("body_en", "")
        if language == "th" and "body_th" in row:
            body = row["body_th"]

        add_text_box(
            slide,
            left=1.2,
            top=top + 0.4,
            width=11.333,
            height=0.8,
            text=body,
            fontsize=SMALL_FONTSIZE,
            font_name=pick_font(theme, "body", language),
            color=text_color,
        )


def add_process_flow(
    slide, theme: Dict[str, Any], content: Dict[str, Any], language: str
) -> None:
    """
    Add a process flow slide with sequential steps and arrows.

    Args:
        slide: Slide object
        theme: Theme dictionary
        content: Slide content with 'title' and 'steps' (list of strings)
        language: Language mode
    """
    colors = theme["colors"]
    bg_color = hex_to_rgb(colors["background"])
    primary_color = hex_to_rgb(colors["primary"])
    accent_color = hex_to_rgb(colors["accent"])
    text_color = hex_to_rgb(colors["text"])

    add_rectangle_shape(slide, 0, 0, 13.333, 7.5, bg_color)

    # Title
    title = content.get("title") or content.get("title_en", "")
    if language == "th" and "title_th" in content:
        title = content["title_th"]

    add_text_box(
        slide,
        left=0.5,
        top=0.5,
        width=12.333,
        height=0.6,
        text=title,
        fontsize=TITLE_FONTSIZE - 8,
        font_name=pick_font(theme, "header", language),
        color=primary_color,
        bold=True,
    )

    # Steps
    steps = content.get("steps", [])
    step_width = 2.0
    step_height = 1.2
    total_width = (step_width + 0.3) * len(steps)
    start_left = (13.333 - total_width) / 2
    start_top = 3.0

    for idx, step in enumerate(steps[:5]):  # Limit to 5 steps
        left = start_left + (idx * (step_width + 0.3))

        # Step box
        add_rectangle_shape(
            slide,
            left=left,
            top=start_top,
            width=step_width,
            height=step_height,
            fill_color=primary_color,
            line_color=accent_color,
        )

        # Step text
        step_text = step if isinstance(step, str) else str(step)
        add_text_box(
            slide,
            left=left + 0.1,
            top=start_top + 0.1,
            width=step_width - 0.2,
            height=step_height - 0.2,
            text=step_text,
            fontsize=SMALL_FONTSIZE,
            font_name=pick_font(theme, "body", language),
            color=bg_color,
            align=PP_ALIGN.CENTER,
        )

        # Arrow (text placeholder if last step)
        if idx < len(steps) - 1:
            add_text_box(
                slide,
                left=left + step_width,
                top=start_top + 0.4,
                width=0.25,
                height=0.4,
                text="→",
                fontsize=20,
                font_name=SYMBOL_FONT,
                color=primary_color,
                align=PP_ALIGN.CENTER,
            )


def add_phased_timeline(
    slide, theme: Dict[str, Any], content: Dict[str, Any], language: str
) -> None:
    """
    Add a phased timeline slide.

    Args:
        slide: Slide object
        theme: Theme dictionary
        content: Slide content with 'title' and 'phases' (list of {name, duration, description})
        language: Language mode
    """
    colors = theme["colors"]
    bg_color = hex_to_rgb(colors["background"])
    primary_color = hex_to_rgb(colors["primary"])
    text_color = hex_to_rgb(colors["text"])
    muted_color = hex_to_rgb(colors["muted"])

    add_rectangle_shape(slide, 0, 0, 13.333, 7.5, bg_color)

    # Title
    title = content.get("title") or content.get("title_en", "")
    if language == "th" and "title_th" in content:
        title = content["title_th"]

    add_text_box(
        slide,
        left=0.5,
        top=0.5,
        width=12.333,
        height=0.6,
        text=title,
        fontsize=TITLE_FONTSIZE - 8,
        font_name=pick_font(theme, "header", language),
        color=primary_color,
        bold=True,
    )

    # Phases
    phases = content.get("phases", [])
    phase_width = 2.5
    start_top = 1.8
    start_left = 0.5

    for idx, phase in enumerate(phases[:4]):  # Limit to 4 phases
        left = start_left + (idx * phase_width)

        # Phase name
        name = phase.get("name") or phase.get("name_en", "")
        if language == "th" and "name_th" in phase:
            name = phase["name_th"]

        add_text_box(
            slide,
            left=left,
            top=start_top,
            width=phase_width - 0.2,
            height=0.5,
            text=name,
            fontsize=BODY_FONTSIZE,
            font_name=pick_font(theme, "header", language),
            color=primary_color,
            bold=True,
        )

        # Duration
        duration = phase.get("duration", "")
        add_text_box(
            slide,
            left=left,
            top=start_top + 0.6,
            width=phase_width - 0.2,
            height=0.4,
            text=duration,
            fontsize=SMALL_FONTSIZE,
            font_name=pick_font(theme, "body", language),
            color=muted_color,
        )

        # Description
        description = phase.get("description") or phase.get("description_en", "")
        if language == "th" and "description_th" in phase:
            description = phase["description_th"]

        add_text_box(
            slide,
            left=left,
            top=start_top + 1.2,
            width=phase_width - 0.2,
            height=4.5,
            text=description,
            fontsize=SMALL_FONTSIZE,
            font_name=pick_font(theme, "body", language),
            color=text_color,
        )


def add_table_slide(
    slide, theme: Dict[str, Any], content: Dict[str, Any], language: str
) -> None:
    """
    Add a table slide.

    Args:
        slide: Slide object
        theme: Theme dictionary
        content: Slide content with 'title', 'headers', and 'rows'
        language: Language mode
    """
    colors = theme["colors"]
    bg_color = hex_to_rgb(colors["background"])
    primary_color = hex_to_rgb(colors["primary"])
    text_color = hex_to_rgb(colors["text"])
    secondary_color = hex_to_rgb(colors["secondary"])

    add_rectangle_shape(slide, 0, 0, 13.333, 7.5, bg_color)

    # Title
    title = content.get("title") or content.get("title_en", "")
    if language == "th" and "title_th" in content:
        title = content["title_th"]

    add_text_box(
        slide,
        left=0.5,
        top=0.5,
        width=12.333,
        height=0.6,
        text=title,
        fontsize=TITLE_FONTSIZE - 8,
        font_name=pick_font(theme, "header", language),
        color=primary_color,
        bold=True,
    )

    # Create table shape
    headers = content.get("headers", [])
    rows = content.get("rows", [])

    if not headers or not rows:
        logger.warning("Table slide has no headers or rows")
        return

    num_cols = len(headers)
    num_rows = min(len(rows) + 1, 6)  # +1 for header row, limit to 6

    left = Inches(0.75)
    top = Inches(1.4)
    width = Inches(11.833)
    height = Inches(5.5)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary_color

        p = cell.text_frame.paragraphs[0]
        p.text = str(header)
        p.font.size = Pt(SMALL_FONTSIZE)
        p.font.bold = True
        p.font.color.rgb = bg_color
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows[:num_rows - 1]):
        for col_idx, cell_value in enumerate(row if isinstance(row, (list, tuple)) else [row]):
            cell = table_shape.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = secondary_color if row_idx % 2 == 0 else bg_color

            p = cell.text_frame.paragraphs[0]
            p.text = str(cell_value)
            p.font.size = Pt(SMALL_FONTSIZE)
            p.font.color.rgb = text_color


def build_deck(
    outline_path: str,
    theme_path: str,
    language: str,
    output_path: str,
    validate_only: bool = False,
    debug: bool = False,
) -> None:
    """
    Build a PPTX deck from outline and theme JSON files.

    Args:
        outline_path: Path to outline JSON file
        theme_path: Path to theme JSON file
        language: Language mode ('th', 'en', or 'bilingual')
        output_path: Path for output PPTX file
        validate_only: If True, validate JSON only without rendering
        debug: If True, enable debug logging

    Raises:
        FileNotFoundError: If input files don't exist
        json.JSONDecodeError: If JSON files are invalid
        ValueError: If layout type is unrecognized
    """
    if debug:
        logger.setLevel(logging.DEBUG)

    logger.info(f"Loading outline from {outline_path}")
    outline_path_obj = Path(outline_path)
    if not outline_path_obj.exists():
        raise FileNotFoundError(f"Outline file not found: {outline_path}")

    with open(outline_path_obj, "r", encoding="utf-8") as f:
        outline = json.load(f)

    logger.info(f"Loading theme from {theme_path}")
    theme_path_obj = Path(theme_path)
    if not theme_path_obj.exists():
        raise FileNotFoundError(f"Theme file not found: {theme_path}")

    with open(theme_path_obj, "r", encoding="utf-8") as f:
        theme = json.load(f)

    logger.info(f"Outline validation: OK")
    logger.info(f"Theme validation: OK")
    logger.info(f"Language mode: {language}")

    if validate_only:
        logger.info("Validate-only mode: JSON files are valid. Exiting without rendering.")
        return

    # Create presentation
    if ICE_TEMPLATE:
        if not Path(ICE_TEMPLATE).is_file():
            raise FileNotFoundError(f"ICE_TEMPLATE ชี้ไปไฟล์ที่ไม่มีอยู่จริง: {ICE_TEMPLATE}")
        prs = Presentation(ICE_TEMPLATE)          # ขนาดสไลด์มาจากแม่แบบ (16:9 13.333 × 7.5 นิ้ว)
        _template_remove_slides(prs)
        logger.info(f"Template mode: {ICE_TEMPLATE} · layouts: {', '.join(l.name for l in prs.slide_layouts)}")
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

    slides_data = outline.get("slides", [])
    total_slides = len(slides_data)

    logger.info(f"Building {total_slides} slides...")

    layout_mapping = {
        "title-hero": add_title_hero,
        "two-column": add_two_column,
        "stat-callout": add_stat_callout,
        "icon-text-rows": add_icon_text_rows,
        "process-flow": add_process_flow,
        "phased-timeline": add_phased_timeline,
        "table": add_table_slide,
    }

    for slide_idx, slide_content in enumerate(slides_data):
        layout_type = slide_content.get("layout")

        if layout_type not in layout_mapping:
            raise ValueError(
                f"Unrecognized layout type '{layout_type}' on slide {slide_idx + 1}. "
                f"Supported types: {', '.join(layout_mapping.keys())}"
            )

        # Add blank slide
        if ICE_TEMPLATE:
            tpl_name = slide_content.get("template_layout") or TEMPLATE_DEFAULT_LAYOUT.get(
                layout_type, TEMPLATE_FALLBACK_LAYOUT)
            tpl_layout = _template_layout(prs, tpl_name)
            slide = prs.slides.add_slide(tpl_layout)
            _template_prepare_slide(slide, tpl_layout)
        else:
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

        # Apply layout function
        layout_fn = layout_mapping[layout_type]
        try:
            layout_fn(slide, theme, slide_content, language)
            logger.debug(f"Slide {slide_idx + 1}: {layout_type} OK")
        except Exception as e:
            logger.error(f"Error rendering slide {slide_idx + 1} ({layout_type}): {e}")
            raise

        # Add footer (โหมดแม่แบบ: footer + เลขหน้ามาจาก layout แล้ว ไม่วาดซ้ำ)
        if not ICE_TEMPLATE:
            add_footer(slide, theme, slide_idx + 1, total_slides)

    # Save
    output_path_obj = Path(output_path)
    output_path_obj.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path_obj)

    # Summary
    file_size_mb = output_path_obj.stat().st_size / (1024 * 1024)
    logger.info(f"Deck built successfully.")
    logger.info(f"Slide count: {total_slides}")
    logger.info(f"File size: {file_size_mb:.2f} MB")
    logger.info(f"Output: {output_path_obj.absolute()}")


def main() -> None:
    """Parse CLI arguments and build deck."""
    parser = argparse.ArgumentParser(
        description="Build a B2B PPTX deck from outline and theme JSON files."
    )
    parser.add_argument(
        "--outline",
        required=True,
        help="Path to outline JSON file",
    )
    parser.add_argument(
        "--theme",
        required=True,
        help="Path to theme JSON file",
    )
    parser.add_argument(
        "--language",
        required=True,
        choices=["th", "en", "bilingual"],
        help="Language mode",
    )
    parser.add_argument(
        "--output",
        required=True,
        help="Path for output PPTX file",
    )
    parser.add_argument(
        "--debug",
        action="store_true",
        help="Enable debug logging",
    )
    parser.add_argument(
        "--validate-only",
        action="store_true",
        help="Validate JSON files without rendering",
    )

    args = parser.parse_args()

    try:
        build_deck(
            outline_path=args.outline,
            theme_path=args.theme,
            language=args.language,
            output_path=args.output,
            validate_only=args.validate_only,
            debug=args.debug,
        )
    except (FileNotFoundError, json.JSONDecodeError, ValueError) as e:
        logger.error(f"Error: {e}")
        sys.exit(1)
    except Exception as e:
        logger.error(f"Unexpected error: {e}", exc_info=args.debug)
        sys.exit(1)


if __name__ == "__main__":
    main()
