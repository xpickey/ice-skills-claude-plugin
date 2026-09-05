#!/usr/bin/env python3
"""
build_pptx.py — generic PPTX builder for sales agents.

Usage:
    python3 build_pptx.py spec.json out.pptx
    ICE_TEMPLATE=/path/iCE-Propose_Master.pptx python3 build_pptx.py spec.json out.pptx   (V02R05 โหมดแม่แบบ)

FONT RULE (V02R04 · 2026.08.05) — ⚠ ข้อความเดิมที่ว่า "default = Tahoma ทุกข้อความ"
  **ยกเลิกแล้ว** (โค้ดไม่เคยอ่าน docstring นี้ · ตั้งแต่ V02R01 ฟอนต์มาจาก font_policy.RAILS)
  • ฟอนต์มาจาก **ราง** ที่ infer_rail() เดาจากชนิดเอกสาร (เอกชน/ราชการ)
  • สไลด์แน่น (>400 ตัวอักษร หรือ >8 บรรทัด หรือ ตาราง >40 ช่อง ในสไลด์ใดสไลด์หนึ่ง)
    → สลับทั้งเด็คเป็น 'Leelawadee' (ยอดวรรณยุกต์เตี้ยสุด = ไม่ชนเมื่อบีบบรรทัด)
    ปิด/บังคับด้วย spec["dense"] = false / true
  • spec["font_family"] ระบุเอง = เคารพ ไม่แทรกแซง (แต่ยังผ่านด่านนโยบาย V1/V2/V4/V5)

spec.json schema:
{
  "title": "Deck title",
  "subtitle": "Optional subtitle",
  "author": "Optional",
  "font_family": "Tahoma",   # optional, default = Tahoma (TH+EN balanced)
  "theme": {"primary": "#1F4E79", "accent": "#2E75B6"},   # optional
  "slides": [
    {"layout": "title", "title": "...", "subtitle": "..."},
    {"layout": "section", "title": "Section header"},
    {"layout": "bullets", "title": "...", "bullets": ["...", "..."]},
    {"layout": "two_column", "title": "...", "left": ["..."], "right": ["..."]},
    {"layout": "table", "title": "...", "headers": ["A","B"], "rows": [["x","y"]]},
    {"layout": "kpi", "title": "...", "kpis": [{"label":"ARR","value":"$1.2M","delta":"+18% YoY"}]},
    {"layout": "image", "title": "...", "image_path": "/path/to.png"},
    {"layout": "thanks", "title": "Thank you", "subtitle": "..."}
  ]
}
"""
import sys, os, json, subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls


def OxmlElement(tag):
    return parse_xml(f'<{tag} {nsdecls("a")}/>')


# ⭐ นโยบายฟอนต์มาจาก SSOT เดียว — ห้าม hard-code ชื่อฟอนต์ในไฟล์นี้ (V02R02)
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from font_policy import (RAILS, resolve_font_policy, infer_rail,
                         measure_slide_density, DENSE_FONT)   # noqa: E402


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── CHAR GUARD (PPTX Lesson #18) ────────────────────────────────────────────
# U+2192 "→" (และ arrow ญาติ) ทำให้ PowerPoint for Mac ปฏิเสธทั้งไฟล์ (Repair)
# ขณะที่ LibreOffice/qlmanage ปล่อยผ่าน (false-green). แทนด้วย ▸ (U+25B8) ที่
# เปิดได้ทุก engine + สื่อความ flow เดียวกัน. ตรวจซ้ำที่ deck_qa.py (safety net).
CHAR_REPLACEMENTS = {
    "→": "▸",  # → RIGHTWARDS ARROW       → ▸ BLACK RIGHT-POINTING SMALL TRIANGLE
    "⟶": "▸",  # ⟶ LONG RIGHTWARDS ARROW  → ▸
    "➜": "▸",  # ➜ HEAVY ROUND-TIPPED ARROW → ▸
    "➔": "▸",  # ➔ HEAVY WIDE-HEADED ARROW  → ▸
    "➙": "▸",  # ➙ HEAVY RIGHTWARDS ARROW   → ▸
}


def _sanitize_chars(obj, _stats=None):
    """Recursively replace PowerPoint-rejecting chars in any string within spec.
    Returns (sanitized_obj, replacement_count). Logs to stderr — never silent."""
    top = _stats is None
    if _stats is None:
        _stats = {"count": 0}
    if isinstance(obj, str):
        out = obj
        for bad, good in CHAR_REPLACEMENTS.items():
            if bad in out:
                _stats["count"] += out.count(bad)
                out = out.replace(bad, good)
        return (out, _stats["count"]) if top else out
    if isinstance(obj, dict):
        res = {k: _sanitize_chars(v, _stats) for k, v in obj.items()}
        return (res, _stats["count"]) if top else res
    if isinstance(obj, list):
        res = [_sanitize_chars(v, _stats) for v in obj]
        return (res, _stats["count"]) if top else res
    return (obj, _stats["count"]) if top else obj


def add_title_slide(prs, slide, theme):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    if len(s.placeholders) > 1:
        s.placeholders[1].text = slide.get("subtitle", "")


def add_section_slide(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")


def add_bullets_slide(prs, slide, theme):
    layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    body = s.placeholders[1].text_frame
    body.text = slide["bullets"][0] if slide.get("bullets") else ""
    for b in slide.get("bullets", [])[1:]:
        p = body.add_paragraph()
        p.text = b


def add_two_column(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    left_items = slide.get("left", [])
    right_items = slide.get("right", [])
    left = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)).text_frame
    right = s.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(4.5), Inches(5)).text_frame
    if left_items:
        left.text = left_items[0]
        for t in left_items[1:]:
            left.add_paragraph().text = t
    if right_items:
        right.text = right_items[0]
        for t in right_items[1:]:
            right.add_paragraph().text = t


def add_table(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    headers = slide.get("headers", [])
    rows = slide.get("rows", [])
    if not headers:
        return
    rows_n = len(rows) + 1
    cols_n = len(headers)
    tbl = s.shapes.add_table(rows_n, cols_n, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * rows_n)).table
    for j, h in enumerate(headers):
        tbl.cell(0, j).text = str(h)
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            tbl.cell(i, j).text = str(v)


def add_kpi(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    kpis = slide.get("kpis", [])
    n = max(1, len(kpis))
    box_w = 9 / n
    for i, k in enumerate(kpis):
        left = Inches(0.5 + i * box_w)
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(box_w - 0.2), Inches(2.5))
        tf = shape.text_frame
        tf.text = k.get("label", "")
        p = tf.add_paragraph()
        p.text = k.get("value", "")
        p.font.size = Pt(28)
        p.font.bold = True
        if k.get("delta"):
            p2 = tf.add_paragraph()
            p2.text = k["delta"]


def add_image(prs, slide, theme):
    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.get("title", "")
    if slide.get("image_path"):
        s.shapes.add_picture(slide["image_path"], Inches(0.5), Inches(1.5), Inches(9), Inches(5))


# ── ⭐ D1 TRI-SLOT FONT BINDING (V02R01 — 2026.08.04) ────────────────────────
# บั๊กที่แก้: ไฟล์นี้ **ไม่เคย set ฟอนต์เลยสักบรรทัด** ตลอดอายุการใช้งาน — docstring
# บอก "default Tahoma" และ spec key `font_family` แต่โค้ดไม่เคยอ่านทั้งคู่ (dead doc)
# → output ได้ฟอนต์ default ของ python-pptx (Calibri) ซึ่งไม่มี glyph ไทยเลย
# ทำเป็น post-pass เดินทุก run แทนการแก้ layout function ทั้ง 8 ตัว (ปลอดภัยกว่า)
def _bind_font(tf, font):
    for p in tf.paragraphs:
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            for tag in ("a:latin", "a:ea", "a:cs"):          # ครบ 3 slot — D1
                for old in rPr.findall(qn(tag)):
                    rPr.remove(old)
                el = OxmlElement(tag)
                el.set("typeface", font)
                rPr.append(el)


def apply_font(prs, font):
    """เดินทุก shape/table cell ในทุก slide แล้วผูกฟอนต์ครบ 3 slot"""
    n = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                _bind_font(sh.text_frame, font); n += 1
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        _bind_font(cell.text_frame, font); n += 1
    return n


# ═══════════════════════════════════════════════════════════════════════════════
# ⭐ V02R05 (2026.09.05 · Wave B) — โหมดแม่แบบ: ICE_TEMPLATE=<path ไฟล์ .pptx แม่แบบ>
#   ต้นตอ: builder ทุกตัวสร้าง deck จาก Presentation() เปล่าแล้ววาดทุกอย่างจาก shape primitive
#   (เฉลี่ย 47 shape ต่อหน้า วัตถุทับกัน ฟอนต์ไม่คงที่ หน้าตาไม่เหมือน CI) — โหมดนี้เปิดแม่แบบจริง
#   ที่ฝังฟอนต์/สี/กริดของ iCE ไว้แล้ว (b2b-slide-designer/assets/masters/iCE-Propose_Master.pptx)
#   แล้ว "เติมเนื้อหาลง placeholder" ของ layout ที่เลือกตามชื่อเท่านั้น
#   ⚠ ไม่ตั้งตัวแปรนี้ = พฤติกรรมเดิมทุกอย่าง 100% (โค้ดด้านบนไม่ถูกแตะ)
#
#   ชื่อ layout ใน spec ใช้ได้ทั้งชื่อเดิม (title/section/bullets/two_column/table/kpi/image/thanks)
#   และชื่อของแม่แบบ (cover/divider/action-title-body/two-column/three-card/table/timeline/closing/appendix)
#   คีย์เพิ่มเติมที่โหมดแม่แบบอ่าน (ไม่บังคับ):
#     ระดับ deck : "footer" (ข้อความท้ายหน้า) · "icon_color" (hex ของ icon ที่แปลงจาก SVG)
#     ทุกหน้าเนื้อหา : "icon" = ชื่อไฟล์ mdi-*.svg ในคลัง หรือ path .png/.svg → icon นำหน้าหัวเรื่อง
#     cover  : "kicker" "subtitle" "meta"          · divider : "number" "subtitle"
#     action-title-body : "bullets" + "image_path" หรือ "image_icon" (พื้นที่ภาพด้านขวา)
#     two-column : "left_title" "left" "left_icon" "right_title" "right" "right_icon"
#     three-card : "cards": [{"icon","title","text"|"bullets"}] (หรือแปลงจาก "kpis" ให้อัตโนมัติ)
#     table  : "headers" "rows" "note"              · timeline : "phases": [{"label","text"}] ≤4
#     closing: "subtitle" "contact"                 · appendix : "bullets"
#   bullet ระดับสอง: ข้อความขึ้นต้นด้วยสองช่องว่าง หรือ {"text": "...", "level": 1}
# ═══════════════════════════════════════════════════════════════════════════════
ICE_TEMPLATE = os.environ.get("ICE_TEMPLATE") or None
ICON_DIR = os.path.expanduser("~/.claude/skills/b2b-slide-designer/assets/icons")

TEMPLATE_LAYOUT_ALIAS = {
    "title": "cover", "cover": "cover",
    "section": "divider", "divider": "divider",
    "bullets": "action-title-body", "image": "action-title-body", "action-title-body": "action-title-body",
    "two_column": "two-column", "two-column": "two-column",
    "kpi": "three-card", "cards": "three-card", "three-card": "three-card",
    "table": "table",
    "timeline": "timeline",
    "thanks": "closing", "closing": "closing",
    "appendix": "appendix",
}


def _layout_by_name(prs, name):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    avail = ", ".join(l.name for l in prs.slide_layouts)
    sys.exit(f"แม่แบบไม่มี layout ชื่อ '{name}' (มี: {avail})")


def _remove_all_slides(prs):
    """ลบสไลด์ตัวอย่างที่ติดมากับแม่แบบ — เหลือแต่ master/layout"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        prs.part.drop_rel(sldId.rId)
        sldIdLst.remove(sldId)


def _ph(slide, idx):
    for shp in slide.placeholders:
        if shp.placeholder_format.idx == idx:
            return shp
    return None


def _set_text(ph, items):
    """เติมข้อความลง placeholder · str = ย่อหน้าเดียว · list = หลายย่อหน้า (bullet ตาม layout)"""
    if ph is None or items in (None, "", []):
        return False
    if isinstance(items, str):
        items = [items]
    tf = ph.text_frame
    first = True
    for it in items:
        level = 0
        if isinstance(it, dict):
            level, it = int(it.get("level", 0)), str(it.get("text", ""))
        elif it.startswith("  "):
            level, it = 1, it.strip()
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = it
        p.level = level
        first = False
    return True


def _icon_png(src, color, out_dir):
    """คืน path PNG ของ icon · รับชื่อ mdi-xxx / ไฟล์ .svg / ไฟล์ .png
    SVG → PNG ใช้ qlmanage ของ macOS (ไม่ต้องติดตั้งอะไร) · ย้อมสีด้วยการแทน currentColor ก่อนแปลง"""
    if not src:
        return None
    if src.lower().endswith(".png") and os.path.isfile(src):
        return src
    name = os.path.basename(src)
    if not name.endswith(".svg"):
        name += ".svg"
    if not name.startswith("mdi-") and not os.path.isfile(src):
        name = "mdi-" + name
    svg_path = src if os.path.isfile(src) else os.path.join(ICON_DIR, name)
    if not os.path.isfile(svg_path):
        sys.exit(f"ไม่พบ icon: {src} (ค้นใน {ICON_DIR} — ดูรายชื่อที่ INDEX.md)")
    icon_dir = os.path.join(out_dir, "_icons")
    os.makedirs(icon_dir, exist_ok=True)
    hexcol = (color or "1E66A4").lstrip("#").upper()
    stem = os.path.splitext(os.path.basename(svg_path))[0]
    png = os.path.join(icon_dir, f"{stem}-{hexcol}.png")
    if os.path.isfile(png):
        return png
    tinted = os.path.join(icon_dir, f"{stem}-{hexcol}.svg")
    with open(svg_path, encoding="utf-8") as f:
        svg = f.read().replace("currentColor", f"#{hexcol}")
    with open(tinted, "w", encoding="utf-8") as f:
        f.write(svg)
    r = subprocess.run(["qlmanage", "-t", "-s", "512", "-o", icon_dir, tinted],
                       capture_output=True, text=True)
    made = tinted + ".png"
    if not os.path.isfile(made):
        sys.exit(f"แปลง SVG เป็น PNG ไม่สำเร็จ ({r.stderr.strip()[:200]}) — ดูวิธีทางเลือกใน "
                 f"{ICON_DIR}/INDEX.md หัวข้อ 'การแปลง SVG เป็น PNG'")
    os.remove(tinted)
    # qlmanage คืนภาพพื้นขาวทึบ → คำนวณความโปร่งใสจากความเข้มของพิกเซล แล้วระบายสีที่ต้องการ
    # (icon ของคลังเป็นสีเดียวบนพื้นขาว จึงใช้ช่องสีต่ำสุดเป็นตัวบอกว่าพิกเซลนั้นเป็นเนื้อ icon แค่ไหน)
    try:
        from PIL import Image
        rr, gg, bb = int(hexcol[0:2], 16), int(hexcol[2:4], 16), int(hexcol[4:6], 16)
        floor = min(rr, gg, bb)
        with Image.open(made) as im:
            im = im.convert("RGBA")
            px = im.load()
            w, h = im.size
            for yy in range(h):
                for xx in range(w):
                    r, g, b, _ = px[xx, yy]
                    a = 255 - min(r, g, b)
                    a = int(round(a * 255 / max(1, 255 - floor)))
                    px[xx, yy] = (rr, gg, bb, max(0, min(255, a)))
            im.save(png)
        os.remove(made)
    except Exception as e:                      # ไม่มี PIL → ใช้ภาพพื้นขาวไปก่อน (แจ้งเสมอ)
        print(f"⚠ icon {stem}: ทำพื้นโปร่งใสไม่ได้ ({e}) — ใช้ภาพพื้นขาว", file=sys.stderr)
        os.replace(made, png)
    return png


def _fill_pic(slide, ph, path):
    """วางภาพให้พอดีในกรอบ placeholder แบบไม่ตัดขอบ (contain) แล้วถอด placeholder ออก"""
    if ph is None or not path:
        return False
    try:
        from PIL import Image
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = 1, 1
    bx, by, bw, bh = ph.left, ph.top, ph.width, ph.height
    scale = min(bw / iw, bh / ih)
    w, h = int(iw * scale), int(ih * scale)
    x, y = bx + (bw - w) // 2, by + (bh - h) // 2
    slide.shapes.add_picture(path, x, y, w, h)
    ph._element.getparent().remove(ph._element)
    return True


def _clone_footer(slide, layout, footer_text):
    """ยก footer (idx 11) และเลขหน้า (idx 12) จาก layout ลงสไลด์ — python-pptx ไม่ clone ให้เอง
    (PowerPoint แสดงสองอย่างนี้เฉพาะเมื่อสไลด์มี placeholder ของตัวเอง)"""
    import copy as _copy
    spTree = slide.shapes._spTree
    for shp in layout.placeholders:
        idx = shp.placeholder_format.idx
        if idx not in (11, 12):
            continue
        el = _copy.deepcopy(shp._element)
        spTree.append(el)
    for shp in slide.placeholders:
        if shp.placeholder_format.idx == 11 and footer_text:
            shp.text_frame.paragraphs[0].text = footer_text


def _prune_empty(slide):
    """ถอด placeholder ที่ไม่ได้เติม (กันข้อความ 'Click to add' และกล่องว่างค้างในไฟล์)"""
    for shp in list(slide.placeholders):
        pf = shp.placeholder_format
        if pf.idx in (11, 12):
            continue
        if shp.has_text_frame and shp.text_frame.text.strip():
            continue
        if not shp.has_text_frame and getattr(shp, "has_table", False) and shp.has_table:
            continue
        shp._element.getparent().remove(shp._element)


def _title_icon(slide, spec_slide, deck, out_dir):
    icon = spec_slide.get("icon")
    if icon:
        _fill_pic(slide, _ph(slide, 10), _icon_png(icon, deck.get("icon_color"), out_dir))


def _tpl_slide(prs, name):
    layout = _layout_by_name(prs, TEMPLATE_LAYOUT_ALIAS.get(name, name))
    return prs.slides.add_slide(layout), layout


def tpl_cover(prs, sl, deck, out_dir):
    s, lay = _tpl_slide(prs, "cover")
    s.shapes.title.text = sl.get("title", deck.get("title", ""))
    _set_text(_ph(s, 1), sl.get("kicker", deck.get("kicker")))
    _set_text(_ph(s, 2), sl.get("subtitle", deck.get("subtitle")))
    _set_text(_ph(s, 3), sl.get("meta"))
    return s, lay


def tpl_divider(prs, sl, deck, out_dir):
    s, lay = _tpl_slide(prs, "divider")
    s.shapes.title.text = sl.get("title", "")
    _set_text(_ph(s, 1), sl.get("number"))
    _set_text(_ph(s, 2), sl.get("subtitle"))
    return s, lay


def tpl_action_title_body(prs, sl, deck, out_dir):
    s, lay = _tpl_slide(prs, "action-title-body")
    s.shapes.title.text = sl.get("title", "")
    _set_text(_ph(s, 1), sl.get("bullets"))
    if sl.get("image_path"):
        _fill_pic(s, _ph(s, 2), sl["image_path"])
    elif sl.get("image_icon"):
        _fill_pic(s, _ph(s, 2), _icon_png(sl["image_icon"], deck.get("icon_color"), out_dir))
    _title_icon(s, sl, deck, out_dir)
    return s, lay


def tpl_two_column(prs, sl, deck, out_dir):
    s, lay = _tpl_slide(prs, "two-column")
    s.shapes.title.text = sl.get("title", "")
    _set_text(_ph(s, 1), sl.get("left_title"))
    _set_text(_ph(s, 2), sl.get("left"))
    _set_text(_ph(s, 3), sl.get("right_title"))
    _set_text(_ph(s, 4), sl.get("right"))
    for idx, key in ((5, "left_icon"), (6, "right_icon")):
        if sl.get(key):
            _fill_pic(s, _ph(s, idx), _icon_png(sl[key], deck.get("icon_color"), out_dir))
    _title_icon(s, sl, deck, out_dir)
    return s, lay


def tpl_three_card(prs, sl, deck, out_dir):
    s, lay = _tpl_slide(prs, "three-card")
    s.shapes.title.text = sl.get("title", "")
    cards = sl.get("cards")
    if not cards and sl.get("kpis"):
        cards = [{"title": k.get("value", ""), "text": [k.get("label", ""), k.get("delta", "")]}
                 for k in sl["kpis"]]
    for i, c in enumerate((cards or [])[:3]):
        base = i * 3
        if c.get("icon"):
            _fill_pic(s, _ph(s, base + 1), _icon_png(c["icon"], deck.get("icon_color"), out_dir))
        _set_text(_ph(s, base + 2), c.get("title"))
        body = c.get("bullets") or c.get("text")
        if isinstance(body, list):
            body = [t for t in body if t]
        _set_text(_ph(s, base + 3), body)
    _title_icon(s, sl, deck, out_dir)
    return s, lay


def tpl_table(prs, sl, deck, out_dir):
    s, lay = _tpl_slide(prs, "table")
    s.shapes.title.text = sl.get("title", "")
    headers, rows = sl.get("headers", []), sl.get("rows", [])
    ph = _ph(s, 1)
    if headers and ph is not None:
        gf = ph.insert_table(len(rows) + 1, len(headers))
        tbl = gf.table
        sz = Pt(16) if len(rows) <= 6 else Pt(14)
        for j, h in enumerate(headers):
            tbl.cell(0, j).text = str(h)
        for i, row in enumerate(rows, start=1):
            for j, v in enumerate(row):
                tbl.cell(i, j).text = str(v)
        for r in tbl.rows:
            r.height = Inches(0.45)
            for c in r.cells:
                for p in c.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = sz
    _set_text(_ph(s, 2), sl.get("note"))
    _title_icon(s, sl, deck, out_dir)
    return s, lay


def tpl_timeline(prs, sl, deck, out_dir):
    s, lay = _tpl_slide(prs, "timeline")
    s.shapes.title.text = sl.get("title", "")
    for i, ph_ in enumerate((sl.get("phases") or [])[:4]):
        _set_text(_ph(s, i + 1), ph_.get("label"))
        _set_text(_ph(s, i + 5), ph_.get("text"))
    _title_icon(s, sl, deck, out_dir)
    return s, lay


def tpl_closing(prs, sl, deck, out_dir):
    s, lay = _tpl_slide(prs, "closing")
    s.shapes.title.text = sl.get("title", "ขอบคุณ")
    _set_text(_ph(s, 1), sl.get("subtitle"))
    _set_text(_ph(s, 2), sl.get("contact"))
    return s, lay


def tpl_appendix(prs, sl, deck, out_dir):
    s, lay = _tpl_slide(prs, "appendix")
    s.shapes.title.text = sl.get("title", "")
    _set_text(_ph(s, 1), sl.get("bullets"))
    _title_icon(s, sl, deck, out_dir)
    return s, lay


TEMPLATE_LAYOUTS = {
    "cover": tpl_cover, "divider": tpl_divider, "action-title-body": tpl_action_title_body,
    "two-column": tpl_two_column, "three-card": tpl_three_card, "table": tpl_table,
    "timeline": tpl_timeline, "closing": tpl_closing, "appendix": tpl_appendix,
}


def build_with_template(prs, spec, out_path):
    """โหมดแม่แบบ — เติม placeholder ตามชื่อ layout · คืนจำนวนสไลด์"""
    out_dir = os.path.dirname(os.path.abspath(out_path))
    _remove_all_slides(prs)
    footer = spec.get("footer", "iCE Consulting · เอกสารลับ")
    for sl in spec["slides"]:
        name = TEMPLATE_LAYOUT_ALIAS.get(sl.get("layout", "bullets"), sl.get("layout"))
        fn = TEMPLATE_LAYOUTS.get(name)
        if fn is None:
            sys.exit(f"layout '{sl.get('layout')}' ไม่มีในแม่แบบ (ใช้ได้: {', '.join(TEMPLATE_LAYOUTS)})")
        s, lay = fn(prs, sl, spec, out_dir)
        _prune_empty(s)
        _clone_footer(s, lay, footer)
    return len(prs.slides)


LAYOUTS = {
    "title": add_title_slide,
    "section": add_section_slide,
    "bullets": add_bullets_slide,
    "two_column": add_two_column,
    "table": add_table,
    "kpi": add_kpi,
    "image": add_image,
    "thanks": add_title_slide,
}


def build(spec_path, out_path):
    with open(spec_path) as f:
        spec = json.load(f)
    # CHAR GUARD (Lesson #18): auto-replace PowerPoint-rejecting chars (→ ▸) before build
    spec, _n = _sanitize_chars(spec)
    if _n:
        print(f"CHAR-GUARD: replaced {_n} arrow char(s) (U+2192/etc → ▸) — would have caused PowerPoint Repair", file=sys.stderr)
    if ICE_TEMPLATE:
        # ⭐ V02R05 โหมดแม่แบบ — ขนาดสไลด์/ฟอนต์ theme/สี/โลโก้/footer มาจากไฟล์แม่แบบทั้งหมด
        if not os.path.isfile(ICE_TEMPLATE):
            sys.exit(f"ICE_TEMPLATE ชี้ไปไฟล์ที่ไม่มีอยู่จริง: {ICE_TEMPLATE}")
        prs = Presentation(ICE_TEMPLATE)
        print(f"🧩 แม่แบบ: {ICE_TEMPLATE} · layout: {', '.join(l.name for l in prs.slide_layouts)}")
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    theme = spec.get("theme", {})
    if not spec.get("slides"):
        spec["slides"] = [{"layout": "title", "title": spec.get("title", "Untitled"), "subtitle": spec.get("subtitle", "")}]
    elif spec["slides"][0].get("layout") not in ("title", "cover"):
        spec["slides"].insert(0, {"layout": "title", "title": spec.get("title", "Untitled"), "subtitle": spec.get("subtitle", "")})
    if ICE_TEMPLATE:
        build_with_template(prs, spec, out_path)
    else:
        for slide in spec["slides"]:
            layout = slide.get("layout", "bullets")
            LAYOUTS.get(layout, add_bullets_slide)(prs, slide, theme)

    # ⭐ D1 — ฟอนต์มาจากราง (§3.0) · spec override ได้เมื่อลูกค้า/แบรนด์บังคับ
    rail, _why = infer_rail(spec, out_path)
    print(f"📄 ราง: {rail}  ({_why})")
    if rail not in RAILS:
        sys.exit(f"rail ต้องเป็น {'|'.join(RAILS)} (ได้: {rail})")
    font = spec.get("font_family") or RAILS[rail]["font"]

    # ⭐ V02R03 (คำสั่ง user 2026.08.05) — สไลด์แน่น/ต้องบีบบรรทัด → DENSE_FONT ทั้งเด็ค
    #   เหตุผลเชิงตัวเลข: ยอดวรรณยุกต์ Leelawadee 0.737 em vs ฟอนต์ราง 0.924 em
    #   → เมื่อบีบ line-height ตัวที่ยอดสูงกว่าชนก่อน · PPTX ฝังฟอนต์ได้จึงไม่ห่วงเครื่องผู้รับ
    #   ปิดได้ด้วย spec["dense"] = false · บังคับเปิดด้วย spec["dense"] = true
    if not spec.get("font_family"):          # ผู้ใช้ระบุฟอนต์เองแล้ว = เคารพ ไม่แทรกแซง
        _forced = spec.get("dense")
        _dense, _dwhy = measure_slide_density(spec.get("slides", []))
        if _forced is True:
            _dense, _dwhy = True, "spec ระบุ dense=true"
        elif _forced is False:
            _dense, _dwhy = False, "spec ระบุ dense=false — ไม่สลับ"
        # ⚠ V02R04 (QA 2026.08.05): งานราชการ — ฟอนต์บังคับของ TOR ชนะกฎความแน่นเสมอ
        #   เจอจาก QA จริง: เด็ค TOR ที่แน่นถูกสลับทิ้งจาก TH Sarabun New → Leelawadee เงียบ ๆ
        #   auto-switch จึงทำเฉพาะราง private · ราง govt ต้องประกาศ dense=true เองเท่านั้น (มีร่องรอย)
        if _dense and rail != "private" and _forced is not True:
            print(f"🎚 สไลด์แน่น ({_dwhy}) แต่เป็นงานราชการ — คงฟอนต์ '{font}' ตามข้อบังคับ · "
                  f"บีบพื้นที่ด้วย line spacing/ลดเนื้อหาแทน · ยืนยันจะสลับจริง → ใส่ \"dense\": true")
            _dense = False
        if _dense and font != DENSE_FONT:
            print(f"🎚 สไลด์แน่น → เปลี่ยนทั้งเด็คเป็น '{DENSE_FONT}': {_dwhy}")
            print(f"   (ยอดวรรณยุกต์ 0.737 em เทียบ '{font}' 0.924 em = ไม่ชนเมื่อบีบบรรทัด"
                  f" · แลกกับไทยเล็กกว่าละตินมากขึ้น · ปิดด้วย \"dense\": false)")
            font = DENSE_FONT
        elif not _dense:
            print(f"🎚 ความแน่น: {_dwhy} → ใช้ฟอนต์ราง")

    # ⭐ V02R02: ด่านนโยบายก่อนสร้างไฟล์ — ผิดนโยบาย = แก้ให้เป็นฟอนต์ราง + แจ้ง (ไม่ fail)
    font, _notices, _ = resolve_font_policy(font, rail, spec)
    for _n in _notices:
        print(_n)
    n = apply_font(prs, font)
    prs.save(out_path)
    print(f"OK: wrote {out_path} with {len(prs.slides)} slides · "
          f"font='{font}' (rail={rail}) ผูกครบ 3 slot ใน {n} text frame")

    # post-build gate — จุดตรวจเดียวกับทุกฟอร์แมต
    subprocess.run([sys.executable, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                                 "audit_fonts.py"), "--rail", rail,
                    *(["--allow-font", font] if spec.get("font_family") else []), out_path])


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("usage: build_pptx.py spec.json out.pptx", file=sys.stderr)
        sys.exit(2)
    build(sys.argv[1], sys.argv[2])
