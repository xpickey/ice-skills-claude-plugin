#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""ตัวตรวจเลย์เอาต์สไลด์อัตโนมัติ — audit_layout.py (V01R01 · 2026.09.05)

ตรวจไฟล์ .pptx ด้วยเครื่องก่อนส่งให้ผู้ตรวจคุณภาพ (อริส) เพื่อให้รอบตรวจของคนใช้กับเนื้อหาและตรรกะเท่านั้น
กฎที่ตรวจมาจาก "แนวทางการทำสไลด์ของ iCE" (b2b-slide-designer/references/pptx-design-doctrine.md):
  ข้อ 1/5  ข้อความต่อหน้าไม่เกินงบคำของโหมด (เอกสารอ่านเอง 75 คำ · นำเสนอสด 25 คำ) → เกิน = เตือน · เกินมาก = ไม่ผ่าน
  ข้อ 2    หน้าเนื้อหาต้องมีภาพหรือ icon อย่างน้อยหนึ่งชิ้น → ไม่มี = ไม่ผ่าน
  ข้อ 5    วลี "วัตถุประสงค์ของหน้า/สไลด์" หรือ "slide objective" ห้ามปรากฏบนสไลด์ → พบ = ไม่ผ่าน (หัวข้อ "Objective & Scope" ของเอกสารไม่นับ)
  เลย์เอาต์ กล่องข้อความล้นขอบสไลด์ = ไม่ผ่าน · กล่องข้อความสองกล่องซ้อนกัน = ไม่ผ่าน
  จำนวนหน้า เกินเพดานของโหมด (เอกสาร 30 · นำเสนอ 20) = เตือน

วิธีใช้:  python3 ~/.claude/agents/_lib/audit_layout.py FILE.pptx [--mode document|presenter] [--json]
ผลลัพธ์:  ตารางรายหน้า + สรุป · exit 0 = ผ่าน/เตือน · exit 2 = ไม่ผ่าน (มีข้อที่ต้องแก้ก่อนส่งตรวจ)
หมายเหตุ: การนับคำภาษาไทยใช้ค่าประมาณ 4.5 ตัวอักษรต่อคำ · ภาพ = รูปภาพหรือกลุ่มรูปทรงที่ไม่มีข้อความ (icon ที่วาด)
"""
import argparse
import json
import re
import sys

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ต้องติดตั้ง python-pptx ก่อน: pip install python-pptx")
    sys.exit(3)

BUDGET = {"document": (75, 120, 30), "presenter": (25, 40, 20)}   # (warn_words, fail_words, max_slides)
THAI = re.compile(r"[฀-๿]+")
LATIN = re.compile(r"[A-Za-z0-9][A-Za-z0-9'’\-./%]*")
OBJECTIVE = re.compile(r"วัตถุประสงค์ของ(หน้า|สไลด์)|objective of (this|the) (slide|page)|(slide|page) objective", re.I)  # หัวข้อ "Objective & Scope" ของเอกสารเป็นเนื้อหาปกติ ไม่นับ


def word_count(text):
    thai = sum(len(m) for m in THAI.findall(text))
    latin = len(LATIN.findall(text))
    return int(round(thai / 4.5)) + latin


def iter_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield sh, True
            for inner, _ in iter_shapes(sh.shapes):
                yield inner, False
        else:
            yield sh, False


def bbox(sh):
    try:
        if sh.left is None or sh.width is None:
            return None
        return (sh.left, sh.top, sh.left + sh.width, sh.top + sh.height)
    except Exception:
        return None


def overlap_ratio(a, b):
    ix = min(a[2], b[2]) - max(a[0], b[0])
    iy = min(a[3], b[3]) - max(a[1], b[1])
    if ix <= 0 or iy <= 0:
        return 0.0
    inter = ix * iy
    small = min((a[2] - a[0]) * (a[3] - a[1]), (b[2] - b[0]) * (b[3] - b[1])) or 1
    return inter / small


def audit_slide(slide, W, H, mode):
    warn_w, fail_w, _ = BUDGET[mode]
    words = 0
    texts = []          # (bbox, text) ของกล่องที่มีข้อความ
    pictures = 0
    drawn = 0           # รูปทรงที่ไม่มีข้อความ (icon/แผนภาพที่วาด) นับเฉพาะที่เล็กกว่า 40% ของหน้า
    overflow = 0
    objective_hit = False
    for sh, is_group in iter_shapes(slide.shapes):
        bb = bbox(sh)
        if bb and not is_group:
            if bb[0] < -0.01 * W or bb[1] < -0.01 * H or bb[2] > 1.01 * W or bb[3] > 1.01 * H:
                overflow += 1
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE or getattr(sh, "image", None) is not None and sh.shape_type != MSO_SHAPE_TYPE.GROUP:
            pictures += 1
            continue
        has_text = getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
        if has_text:
            t = sh.text_frame.text.strip()
            words += word_count(t)
            if OBJECTIVE.search(t):
                objective_hit = True
            if bb:
                texts.append((bb, t[:40]))
        elif bb and not is_group:
            area = (bb[2] - bb[0]) * (bb[3] - bb[1])
            if 0 < area < 0.4 * W * H:
                drawn += 1
    overlaps = []
    for i in range(len(texts)):
        for j in range(i + 1, len(texts)):
            r = overlap_ratio(texts[i][0], texts[j][0])
            if r > 0.3:
                overlaps.append((texts[i][1], texts[j][1], round(r, 2)))
    has_visual = pictures > 0 or drawn >= 3
    is_content = words > 12
    issues = []
    if overflow:
        issues.append(f"ล้นขอบ {overflow} ชิ้น")
    if overlaps:
        issues.append(f"ข้อความซ้อนกัน {len(overlaps)} คู่")
    if is_content and not has_visual:
        issues.append("หน้าเนื้อหาไม่มีภาพหรือ icon")
    if words > fail_w:
        issues.append(f"ข้อความ {words} คำ เกินงบ {warn_w} มาก ต้องแตกเป็น bullet หรือแยกหน้า")
    elif words > warn_w:
        issues.append(f"ข้อความ {words} คำ เกินงบ {warn_w} (เตือน)")
    if objective_hit:
        issues.append("มีคำว่าวัตถุประสงค์ของหน้าบนสไลด์ (ต้องอยู่ใน spec เท่านั้น)")
    fail = bool(overflow or overlaps or (is_content and not has_visual) or words > fail_w or objective_hit)
    return {"words": words, "pictures": pictures, "drawn": drawn, "overflow": overflow, "overlaps": len(overlaps),
            "content": is_content, "visual": has_visual, "issues": issues, "fail": fail}


def audit(path, mode):
    prs = Presentation(path)
    W, H = prs.slide_width, prs.slide_height
    rows = [audit_slide(s, W, H, mode) for s in prs.slides]
    _, _, max_slides = BUDGET[mode]
    n = len(rows)
    fails = [i + 1 for i, r in enumerate(rows) if r["fail"]]
    warns = [i + 1 for i, r in enumerate(rows) if r["issues"] and not r["fail"]]
    summary = {
        "file": path, "mode": mode, "slides": n,
        "avg_words": round(sum(r["words"] for r in rows) / n, 1) if n else 0,
        "max_words": max((r["words"] for r in rows), default=0),
        "content_without_visual": sum(1 for r in rows if r["content"] and not r["visual"]),
        "overflow_slides": sum(1 for r in rows if r["overflow"]),
        "overlap_slides": sum(1 for r in rows if r["overlaps"]),
        "too_many_slides": n > max_slides,
        "fail_slides": fails, "warn_slides": warns,
        "verdict": "FAIL" if fails else ("WARN" if warns or n > max_slides else "PASS"),
    }
    return summary, rows


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("files", nargs="+")
    ap.add_argument("--mode", choices=list(BUDGET), default="document")
    ap.add_argument("--json", action="store_true")
    a = ap.parse_args()
    worst = 0
    for f in a.files:
        s, rows = audit(f, a.mode)
        if a.json:
            print(json.dumps({"summary": s, "slides": rows}, ensure_ascii=False))
        else:
            print(f"== {f} · โหมด {a.mode} · {s['slides']} หน้า · เฉลี่ย {s['avg_words']} คำ/หน้า · สูงสุด {s['max_words']} คำ ==")
            for i, r in enumerate(rows, 1):
                if r["issues"]:
                    print(f"  หน้า {i:2d} [{'ไม่ผ่าน' if r['fail'] else 'เตือน'}] " + " · ".join(r["issues"]))
            if s["too_many_slides"]:
                print(f"  [เตือน] จำนวนหน้า {s['slides']} เกินเพดานของโหมด {BUDGET[a.mode][2]}")
            print(f"  ผล: {s['verdict']} — ไม่ผ่าน {len(s['fail_slides'])} หน้า · เตือน {len(s['warn_slides'])} หน้า")
        worst = max(worst, 2 if s["verdict"] == "FAIL" else 0)
    sys.exit(worst)


if __name__ == "__main__":
    main()
